import os
import sys
import asyncio
import logging
from pathlib import Path
from typing import List, Optional, Any, Dict
from contextlib import asynccontextmanager
import json
from datetime import datetime, timedelta
import traceback
import io
import tempfile
import re
import difflib
from urllib.parse import quote, urlparse
from PIL import Image
from pdf2image import convert_from_path
import pytesseract
import docx
import openpyxl
from pptx import Presentation
import faiss
import numpy as np
import tiktoken
import fitz  # PyMuPDF
import torch
from transformers import AutoTokenizer, AutoModel
from fastapi import FastAPI, HTTPException, UploadFile, File, BackgroundTasks, Body, Form
from fastapi.middleware.cors import CORSMiddleware
from fastapi.responses import JSONResponse, StreamingResponse
from pydantic import BaseModel, Field, HttpUrl
import httpx
from google.cloud import storage
from google.api_core.exceptions import GoogleAPIError, NotFound
import jsonlines
from trafilatura import extract
from trafilatura.settings import use_config
from bs4 import BeautifulSoup
from uuid import uuid4
import uuid
import aiofiles
from enum import Enum
from io import BytesIO
import PyPDF2
from datetime import datetime, timezone
from collections import defaultdict, Counter
import openai

# --- Logging Setup ---
logging.basicConfig(
    level=logging.INFO,
    format='%(asctime)s - %(name)s - %(levelname)s - %(message)s',
    stream=sys.stdout
)
logger = logging.getLogger(__name__)

# --- Configuration ---
GCS_BUCKET_NAME = os.getenv("GCS_BUCKET_NAME", "linkedin-bot-documents")
VECTOR_DB_DIR = Path("./vector_db")

# Ensure necessary local cache directory exists at startup
VECTOR_DB_DIR.mkdir(exist_ok=True)

# Define local and GCS paths for the vector store
FAISS_INDEX_PATH = VECTOR_DB_DIR / "faiss_index.bin"
METADATA_PATH = VECTOR_DB_DIR / "metadata.jsonl"
GCS_INDEX_BLOB_NAME = "vector_store/faiss_index.bin"
GCS_METADATA_BLOB_NAME = "vector_store/metadata.jsonl"
GCS_SETTINGS_BLOB_NAME = "config/settings.json"
GCS_DOCS_PREFIX = "documents/"

CHUNK_SIZE = 300
OVERLAP_SIZE = 50
TOP_K = 6
EMBEDDING_MODEL = "intfloat/e5-large-v2"
GROQ_MODEL = "llama-3.3-70b-versatile"
GROQ_API_URL = "https://api.groq.com/openai/v1/chat/completions"
ALLOWED_EXTENSIONS = {".pdf", ".docx", ".xlsx", ".xls", ".pptx", ".ppt", ".txt"}

# --- API Keys & Global State ---
GROQ_API_KEY = os.getenv("GROQ_API_KEY", "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").strip()
OPENAI_API_KEY = os.getenv("OPENAI_API_KEY", "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").strip()

if not GROQ_API_KEY:
    logger.warning("GROQ_API_KEY environment variable not set. Chat completions will fail.")
logger.info(f"GROQ_API_KEY starts with: {GROQ_API_KEY[:6] if GROQ_API_KEY else 'NOT SET'}")

tokenizer = tiktoken.get_encoding("cl100k_base")
faiss_index: Optional[faiss.Index] = None
chunks: List[str] = []
metadata: List[dict] = []
embedding_tokenizer = None
embedding_model = None
device = torch.device("cuda" if torch.cuda.is_available() else "cpu")
gcs_client = None
openai_client = None

# --- Google Cloud Storage Client ---
def get_gcs_client() -> storage.Client:
    global gcs_client
    if gcs_client is None:
        try:
            gcs_client = storage.Client()
            logger.info("Google Cloud Storage client initialized successfully")
        except Exception as e:
            logger.error(f"Failed to initialize GCS client: {e}")
            raise
    return gcs_client

# --- Embedding Model Initialization ---
def initialize_embedding_model():
    global embedding_tokenizer, embedding_model
    try:
        embedding_tokenizer = AutoTokenizer.from_pretrained(EMBEDDING_MODEL)
        embedding_model = AutoModel.from_pretrained(EMBEDDING_MODEL).to(device)
        logger.info(f"Embedding model {EMBEDDING_MODEL} loaded successfully on {device}")
    except Exception as e:
        logger.error(f"Failed to load embedding model: {e}")
        raise

def average_pool(last_hidden_states, attention_mask):
    last_hidden = last_hidden_states.masked_fill(~attention_mask[..., None].bool(), 0.0)
    return last_hidden.sum(dim=1) / attention_mask.sum(dim=1)[..., None]

async def get_embeddings(texts: List[str]) -> np.ndarray:
    if not embedding_model or not embedding_tokenizer:
        raise RuntimeError("Embedding model not initialized")
    
    def compute_embeddings_sync(model, batch):
        with torch.no_grad():
            inputs = embedding_tokenizer(batch, padding=True, truncation=True, return_tensors="pt").to(device)
            outputs = model(**inputs)
            embeddings = average_pool(outputs.last_hidden_state, inputs['attention_mask'])
            return embeddings.cpu().numpy()
    
    loop = asyncio.get_event_loop()
    embeddings = await loop.run_in_executor(None, compute_embeddings_sync, embedding_model, texts)
    return embeddings

# --- Vector Database Management ---
def initialize_index():
    global faiss_index
    if faiss_index is None:
        dimension = 1024  # E5-large-v2 dimension
        faiss_index = faiss.IndexFlatIP(dimension)
        logger.info("FAISS index initialized")

def save_vector_db_to_gcs():
    try:
        if faiss_index is None:
            logger.warning("No FAISS index to save")
            return
        
        client = get_gcs_client()
        bucket = client.bucket(GCS_BUCKET_NAME)
        
        # Save FAISS index
        index_bytes = faiss.serialize_index(faiss_index).tobytes()
        index_blob = bucket.blob(GCS_INDEX_BLOB_NAME)
        index_blob.upload_from_string(index_bytes, content_type='application/octet-stream')
        
        # Save metadata
        metadata_blob = bucket.blob(GCS_METADATA_BLOB_NAME)
        with metadata_blob.open('w') as f:
            for item in metadata:
                f.write(json.dumps(item) + '\n')
        
        logger.info("Vector database saved to GCS successfully")
    except Exception as e:
        logger.error(f"Failed to save vector database to GCS: {e}")

def load_vector_db_from_gcs():
    global faiss_index, chunks, metadata
    try:
        client = get_gcs_client()
        bucket = client.bucket(GCS_BUCKET_NAME)
        
        # Load FAISS index
        index_blob = bucket.blob(GCS_INDEX_BLOB_NAME)
        if index_blob.exists():
            index_bytes = index_blob.download_as_bytes()
            faiss_index = faiss.deserialize_index(index_bytes)
            logger.info("FAISS index loaded from GCS")
        
        # Load metadata
        metadata_blob = bucket.blob(GCS_METADATA_BLOB_NAME)
        if metadata_blob.exists():
            metadata = []
            chunks = []
            with metadata_blob.open('r') as f:
                for line in f:
                    item = json.loads(line.strip())
                    metadata.append(item)
                    chunks.append(item['text'])
            logger.info(f"Loaded {len(metadata)} documents from GCS")
        
    except Exception as e:
        logger.error(f"Failed to load vector database from GCS: {e}")

# --- Application Lifespan ---
@asynccontextmanager
async def lifespan(app: FastAPI):
    # Startup
    logger.info("Starting LinkedIn Content Creator API...")
    
    try:
        # Initialize embedding model
        initialize_embedding_model()
        
        # Initialize FAISS index
        initialize_index()
        
        # Load existing vector database
        load_vector_db_from_gcs()
        
        # Initialize OpenAI client
        global openai_client
        openai_client = openai.OpenAI(api_key=OPENAI_API_KEY)
        
        logger.info("LinkedIn Content Creator API started successfully")
    except Exception as e:
        logger.error(f"Failed to start API: {e}")
        raise
    
    yield
    
    # Shutdown
    logger.info("Shutting down LinkedIn Content Creator API...")
    save_vector_db_to_gcs()

# --- FastAPI App ---
app = FastAPI(
    title="LinkedIn Content Creator API",
    description="API for document management and content generation for LinkedIn posts",
    version="1.0.0",
    lifespan=lifespan
)

# CORS middleware
app.add_middleware(
    CORSMiddleware,
    allow_origins=["*"],
    allow_credentials=True,
    allow_methods=["*"],
    allow_headers=["*"],
)

# --- Pydantic Models ---
class AskRequest(BaseModel):
    question: str = Field(..., min_length=1, max_length=1000)

class UploadResponse(BaseModel):
    message: str
    filename: str
    chunks_added: int
    total_chunks: int

class AskResponse(BaseModel):
    answer: str
    question: str
    sources: List[dict]
    client_followups: List[str] = Field(default_factory=list)
    sales_followups: List[str] = Field(default_factory=list)

class CrawlRequest(BaseModel):
    url: HttpUrl

class CrawlResponse(BaseModel):
    url: HttpUrl
    title: str | None = None
    author: str | None = None
    date: str | None = None
    hostname: str | None = None
    main_text: str
    comments: str | None = None
    raw_text: str | None = None
    source: str | None = None

class ContentGenerationRequest(BaseModel):
    prompt: str
    creator_style: str
    research_context: str = ""
    max_tokens: int = 1000

class ContentGenerationResponse(BaseModel):
    linkedin_posts: List[str]
    video_scripts: List[str]
    hashtags: List[str]
    engagement_tips: List[str]

# --- Utility Functions ---
def extract_text_from_file(content: bytes, filename: str) -> str:
    """Extract text from various file types"""
    try:
        if filename.lower().endswith('.pdf'):
            # Use PyMuPDF for better text extraction
            doc = fitz.open(stream=content, filetype="pdf")
            text = ""
            for page in doc:
                text += page.get_text()
            doc.close()
            return text
        
        elif filename.lower().endswith('.docx'):
            doc = docx.Document(BytesIO(content))
            return "\n".join([paragraph.text for paragraph in doc.paragraphs])
        
        elif filename.lower().endswith('.txt'):
            return content.decode('utf-8')
        
        elif filename.lower().endswith(('.xlsx', '.xls')):
            workbook = openpyxl.load_workbook(BytesIO(content))
            text = ""
            for sheet in workbook.worksheets:
                for row in sheet.iter_rows(values_only=True):
                    text += " ".join([str(cell) for cell in row if cell]) + "\n"
            return text
        
        elif filename.lower().endswith(('.pptx', '.ppt')):
            prs = Presentation(BytesIO(content))
            text = ""
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        text += shape.text + "\n"
            return text
        
        else:
            return content.decode('utf-8', errors='ignore')
    
    except Exception as e:
        logger.error(f"Error extracting text from {filename}: {e}")
        return f"Error extracting text from {filename}: {str(e)}"

def chunk_text(text: str) -> List[str]:
    """Split text into chunks for vector search"""
    if not text.strip():
        return []
    
    # Use tiktoken for proper tokenization
    tokens = tokenizer.encode(text)
    chunks = []
    
    for i in range(0, len(tokens), CHUNK_SIZE):
        chunk_tokens = tokens[i:i + CHUNK_SIZE + OVERLAP_SIZE]
        chunk_text = tokenizer.decode(chunk_tokens)
        if chunk_text.strip():
            chunks.append(chunk_text.strip())
    
    return chunks

async def query_groq(messages: List[dict]) -> str:
    """Query Groq API for content generation"""
    try:
        async with httpx.AsyncClient() as client:
            response = await client.post(
                GROQ_API_URL,
                headers={
                    "Authorization": f"Bearer {GROQ_API_KEY}",
                    "Content-Type": "application/json"
                },
                json={
                    "model": GROQ_MODEL,
                    "messages": messages,
                    "temperature": 0.7,
                    "max_tokens": 2000
                },
                timeout=30.0
            )
            response.raise_for_status()
            return response.json()["choices"][0]["message"]["content"]
    except Exception as e:
        logger.error(f"Groq API error: {e}")
        raise HTTPException(status_code=500, detail=f"Content generation failed: {str(e)}")

# --- API Endpoints ---
@app.get("/", summary="Health Check", tags=["System"])
def read_root():
    return {
        "status": "healthy",
        "message": "LinkedIn Content Creator API is online and ready",
        "timestamp": datetime.now().isoformat(),
        "version": "1.0.0"
    }

@app.post("/upload", response_model=UploadResponse, summary="Upload and Process a Document", tags=["Document Management"])
async def upload_file(file: UploadFile = File(...)):
    """Upload a document and process it for content generation"""
    try:
        # Validate file type
        file_extension = Path(file.filename).suffix.lower()
        if file_extension not in ALLOWED_EXTENSIONS:
            raise HTTPException(status_code=400, detail=f"File type {file_extension} not supported")
        
        # Read file content
        content = await file.read()
        
        # Extract text
        text = extract_text_from_file(content, file.filename)
        if not text.strip():
            raise HTTPException(status_code=400, detail="No text content found in file")
        
        # Chunk text
        text_chunks = chunk_text(text)
        if not text_chunks:
            raise HTTPException(status_code=400, detail="Failed to process document content")
        
        # Generate embeddings
        embeddings = await get_embeddings(text_chunks)
        
        # Add to vector database
        global faiss_index, chunks, metadata
        
        # Add to FAISS index
        if faiss_index is None:
            initialize_index()
        
        faiss_index.add(embeddings.astype('float32'))
        
        # Add to metadata
        for i, chunk in enumerate(text_chunks):
            metadata.append({
                'id': str(uuid4()),
                'text': chunk,
                'filename': file.filename,
                'chunk_index': i,
                'upload_date': datetime.now().isoformat(),
                'file_size': len(content),
                'file_type': file.content_type
            })
            chunks.append(chunk)
        
        # Save to GCS
        save_vector_db_to_gcs()
        
        return UploadResponse(
            message="Document uploaded and processed successfully",
            filename=file.filename,
            chunks_added=len(text_chunks),
            total_chunks=len(chunks)
        )
    
    except Exception as e:
        logger.error(f"Upload error: {e}")
        raise HTTPException(status_code=500, detail=str(e))

@app.post("/ask", response_model=AskResponse, summary="Ask a Question and get relevant information", tags=["Document Management"])
async def ask_question(request: AskRequest):
    """Ask a question and get relevant information from uploaded documents"""
    try:
        if not faiss_index or not chunks:
            return AskResponse(
                answer="No documents have been uploaded yet. Please upload some documents first.",
                question=request.question,
                sources=[],
                client_followups=[],
                sales_followups=[]
            )
        
        # Generate query embedding
        query_embedding = await get_embeddings([request.question])
        
        # Search similar chunks
        scores, indices = faiss_index.search(query_embedding.astype('float32'), TOP_K)
        
        # Get relevant chunks
        relevant_chunks = []
        sources = []
        
        for score, idx in zip(scores[0], indices[0]):
            if idx < len(chunks) and score > 0.1:  # Similarity threshold
                chunk = chunks[idx]
                meta = metadata[idx] if idx < len(metadata) else {}
                relevant_chunks.append(chunk)
                sources.append({
                    'filename': meta.get('filename', 'Unknown'),
                    'chunk_index': meta.get('chunk_index', 0),
                    'similarity_score': float(score)
                })
        
        if not relevant_chunks:
            return AskResponse(
                answer="I couldn't find any relevant information in the uploaded documents.",
                question=request.question,
                sources=[],
                client_followups=[],
                sales_followups=[]
            )
        
        # Generate answer using Groq
        context = "\n\n".join(relevant_chunks[:3])  # Use top 3 chunks
        prompt = f"""Based on the following context from uploaded documents, answer the question.

Context:
{context}

Question: {request.question}

Please provide a comprehensive answer based on the context provided. If the context doesn't contain enough information to answer the question, say so.

Answer:"""

        messages = [
            {"role": "system", "content": "You are a helpful assistant that answers questions based on uploaded documents."},
            {"role": "user", "content": prompt}
        ]
        
        answer = await query_groq(messages)
        
        # Generate follow-up questions
        followup_prompt = f"""Based on the question "{request.question}" and the answer provided, generate 3 relevant follow-up questions that would help explore this topic further.

Answer: {answer}

Generate 3 follow-up questions:"""

        followup_messages = [
            {"role": "system", "content": "You are a helpful assistant that generates relevant follow-up questions."},
            {"role": "user", "content": followup_prompt}
        ]
        
        followups_response = await query_groq(followup_messages)
        followups = [q.strip() for q in followups_response.split('\n') if q.strip() and '?' in q][:3]
        
        return AskResponse(
            answer=answer,
            question=request.question,
            sources=sources,
            client_followups=followups,
            sales_followups=[]
        )
    
    except Exception as e:
        logger.error(f"Ask error: {e}")
        raise HTTPException(status_code=500, detail=str(e))

@app.post("/generate-content", response_model=ContentGenerationResponse, summary="Generate LinkedIn content", tags=["Content Generation"])
async def generate_content(request: ContentGenerationRequest):
    """Generate LinkedIn posts and video scripts based on research papers and creator style"""
    try:
        # Create the prompt for content generation
        system_prompt = f"""You are a LinkedIn content creator that generates engaging posts and video scripts in the style of {request.creator_style}.

Creator Style: {request.creator_style}

Research Context:
{request.research_context}

Generate content that:
1. Incorporates insights from the research context
2. Matches the creator's style and tone
3. Is engaging and shareable
4. Includes relevant hashtags and engagement tips

Return your response in JSON format with:
- linkedin_posts: Array of 2 different LinkedIn post variations
- video_scripts: Array of 2 different video script variations  
- hashtags: Array of relevant hashtags
- engagement_tips: Array of tips for better engagement"""

        user_prompt = f"Generate LinkedIn content about: {request.prompt}"

        messages = [
            {"role": "system", "content": system_prompt},
            {"role": "user", "content": user_prompt}
        ]
        
        response = await query_groq(messages)
        
        # Parse JSON response
        try:
            content_data = json.loads(response)
            return ContentGenerationResponse(**content_data)
        except json.JSONDecodeError:
            # Fallback if JSON parsing fails
            return ContentGenerationResponse(
                linkedin_posts=[f"Generated post about: {request.prompt}"],
                video_scripts=[f"Generated video script about: {request.prompt}"],
                hashtags=["#linkedin", "#content", "#research"],
                engagement_tips=["Ask questions to engage your audience", "Use relevant hashtags", "Share personal insights"]
            )
    
    except Exception as e:
        logger.error(f"Content generation error: {e}")
        raise HTTPException(status_code=500, detail=str(e))

@app.get("/documents", summary="List all indexed documents and metadata", tags=["Document Management"])
def list_documents():
    """Get list of all uploaded documents"""
    try:
        documents = {}
        for meta in metadata:
            filename = meta.get('filename', 'Unknown')
            if filename not in documents:
                documents[filename] = {
                    'filename': filename,
                    'upload_date': meta.get('upload_date'),
                    'file_size': meta.get('file_size', 0),
                    'file_type': meta.get('file_type', 'Unknown'),
                    'total_chunks': 0
                }
            documents[filename]['total_chunks'] += 1
        
        return list(documents.values())
    except Exception as e:
        logger.error(f"List documents error: {e}")
        raise HTTPException(status_code=500, detail=str(e))

@app.get("/chunks", summary="Get all chunks and their metadata", tags=["Document Management"])
def get_all_chunks():
    """Get all text chunks with metadata"""
    try:
        return [
            {
                'id': meta.get('id', str(i)),
                'text': chunk,
                'filename': meta.get('filename', 'Unknown'),
                'chunk_index': meta.get('chunk_index', i),
                'upload_date': meta.get('upload_date')
            }
            for i, (chunk, meta) in enumerate(zip(chunks, metadata))
        ]
    except Exception as e:
        logger.error(f"Get chunks error: {e}")
        raise HTTPException(status_code=500, detail=str(e))

@app.delete("/documents/{filename}", summary="Delete a Document and its Embeddings", tags=["Document Management"])
async def delete_document(filename: str, background_tasks: BackgroundTasks):
    """Delete a document and all its associated chunks and embeddings"""
    try:
        global faiss_index, chunks, metadata
        
        # Find chunks for this document
        indices_to_remove = []
        for i, meta in enumerate(metadata):
            if meta.get('filename') == filename:
                indices_to_remove.append(i)
        
        if not indices_to_remove:
            raise HTTPException(status_code=404, detail="Document not found")
        
        # Remove from metadata and chunks
        for i in reversed(indices_to_remove):
            del metadata[i]
            del chunks[i]
        
        # Rebuild FAISS index
        if chunks:
            embeddings = await get_embeddings(chunks)
            faiss_index = faiss.IndexFlatIP(embeddings.shape[1])
            faiss_index.add(embeddings.astype('float32'))
        else:
            faiss_index = None
        
        # Save to GCS
        save_vector_db_to_gcs()
        
        return {"message": f"Document {filename} deleted successfully"}
    
    except Exception as e:
        logger.error(f"Delete document error: {e}")
        raise HTTPException(status_code=500, detail=str(e))

@app.post("/crawl", response_model=CrawlResponse, summary="Extract clean content from a URL and index it", tags=["Website"])
async def crawl_and_extract(request: CrawlRequest = Body(...)):
    """Crawl a website and extract content for indexing"""
    try:
        async with httpx.AsyncClient() as client:
            response = await client.get(str(request.url), timeout=30.0)
            response.raise_for_status()
            html_content = response.text
        
        # Extract text using trafilatura
        config = use_config()
        config.set("DEFAULT_EXTRACTION_TIMEOUT", 30)
        
        extracted_text = extract(html_content, config=config)
        if not extracted_text:
            # Fallback to BeautifulSoup
            soup = BeautifulSoup(html_content, 'html.parser')
            extracted_text = soup.get_text()
        
        if not extracted_text.strip():
            raise HTTPException(status_code=400, detail="No content could be extracted from the URL")
        
        # Extract metadata
        soup = BeautifulSoup(html_content, 'html.parser')
        title = soup.find('title')
        title_text = title.get_text() if title else None
        
        # Process the extracted text
        text_chunks = chunk_text(extracted_text)
        if not text_chunks:
            raise HTTPException(status_code=400, detail="Failed to process extracted content")
        
        # Generate embeddings
        embeddings = await get_embeddings(text_chunks)
        
        # Add to vector database
        global faiss_index, chunks, metadata
        
        if faiss_index is None:
            initialize_index()
        
        faiss_index.add(embeddings.astype('float32'))
        
        # Add to metadata
        for i, chunk in enumerate(text_chunks):
            metadata.append({
                'id': str(uuid4()),
                'text': chunk,
                'filename': f"Web Content - {request.url}",
                'chunk_index': i,
                'upload_date': datetime.now().isoformat(),
                'file_size': len(extracted_text),
                'file_type': 'text/html',
                'url': str(request.url)
            })
            chunks.append(chunk)
        
        # Save to GCS
        save_vector_db_to_gcs()
        
        return CrawlResponse(
            url=request.url,
            title=title_text,
            main_text=extracted_text[:1000] + "..." if len(extracted_text) > 1000 else extracted_text,
            source=str(request.url)
        )
    
    except Exception as e:
        logger.error(f"Crawl error: {e}")
        raise HTTPException(status_code=500, detail=str(e))

@app.post("/refresh", summary="Synchronize the index with GCS documents", tags=["Document Management"])
async def refresh_index():
    """Refresh the vector index from Google Cloud Storage"""
    try:
        load_vector_db_from_gcs()
        return {"message": "Index refreshed successfully"}
    except Exception as e:
        logger.error(f"Refresh error: {e}")
        raise HTTPException(status_code=500, detail=str(e))

if __name__ == "__main__":
    import uvicorn
    uvicorn.run(app, host="0.0.0.0", port=8000) 