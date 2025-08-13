import os
import sys
import asyncio
import logging
from pathlib import Path
from typing import List, Optional, Any, Dict
from contextlib import asynccontextmanager
import json
from datetime import datetime, timedelta
import traceback
import io
import tempfile
import re
import difflib
from urllib.parse import quote, urlparse
from PIL import Image
from pdf2image import convert_from_path
import pytesseract
import docx
import openpyxl
from pptx import Presentation
import faiss
import numpy as np
import tiktoken
import fitz  # PyMuPDF
import torch
from transformers import AutoTokenizer, AutoModel
from fastapi import FastAPI, HTTPException, UploadFile, File, BackgroundTasks, Body, Form
from fastapi.middleware.cors import CORSMiddleware
from fastapi.responses import JSONResponse, StreamingResponse
from pydantic import BaseModel, Field, HttpUrl
import httpx
from google.cloud import storage
from google.api_core.exceptions import GoogleAPIError, NotFound
import jsonlines
from trafilatura import extract
from trafilatura.settings import use_config
from bs4 import BeautifulSoup
from uuid import uuid4
import uuid
import aiofiles
from enum import Enum
from io import BytesIO
import PyPDF2
from datetime import datetime, timezone
from collections import defaultdict, Counter
import openai

# --- Logging Setup ---
# Configure logging to be more verbose, especially for startup.
logging.basicConfig(
    level=logging.INFO,
    format='%(asctime)s - %(name)s - %(levelname)s - %(message)s',
    stream=sys.stdout
)
logger = logging.getLogger(__name__)

# --- Configuration ---
GCS_BUCKET_NAME = os.getenv("GCS_BUCKET_NAME", "spikedai-bucket-one")
VECTOR_DB_DIR = Path("./vector_db")

# Ensure necessary local cache directory exists at startup
VECTOR_DB_DIR.mkdir(exist_ok=True)

# Define local and GCS paths for the vector store
FAISS_INDEX_PATH = VECTOR_DB_DIR / "faiss_index.bin"
METADATA_PATH = VECTOR_DB_DIR / "metadata.jsonl"
GCS_INDEX_BLOB_NAME = "vector_store/faiss_index.bin"
GCS_METADATA_BLOB_NAME = "vector_store/metadata.jsonl"
GCS_SETTINGS_BLOB_NAME = "config/settings.json"
GCS_DOCS_PREFIX = "documents/"
BASE_URL = "https://sales-assistant-service-822359826336.asia-south1.run.app"

CHUNK_SIZE = 300
OVERLAP_SIZE = 50
TOP_K = 6
EMBEDDING_MODEL = "intfloat/e5-large-v2"
GROQ_MODEL = "llama-3.3-70b-versatile"
GROQ_API_URL = "https://api.groq.com/openai/v1/chat/completions"
ALLOWED_EXTENSIONS = {".pdf", ".docx", ".xlsx", ".xls", ".pptx", ".ppt"}

# --- API Keys & Global State ---
GROQ_API_KEY = os.getenv("GROQ_API_KEY", "").strip()
VEXA_API_BASE_URL = "https://gateway.dev.vexa.ai"
VEXA_API_KEY = os.getenv("VEXA_API_KEY", "").strip()
OPENAI_API_KEY = os.getenv("OPENAI_API_KEY", "").strip()

if not GROQ_API_KEY:
    logger.warning("GROQ_API_KEY environment variable not set. Chat completions will fail.")
logger.info(f"GROQ_API_KEY starts with: {GROQ_API_KEY[:6] if GROQ_API_KEY else 'NOT SET'}")

tokenizer = tiktoken.get_encoding("cl100k_base")
faiss_index: Optional[faiss.Index] = None
chunks: List[str] = []
metadata: List[dict] = []
embedding_tokenizer = None
embedding_model = None
device = torch.device("cuda" if torch.cuda.is_available() else "cpu")
gcs_client = None
openai_client = None
app_settings: Dict[str, Any] = {
    "botName": "SpikedAI",
    "selectedPersona": "balanced",
    "selectedAnswerStyles": [],
    "customPrompt": ""
} 

# --- GCS Client Singleton ---
def get_gcs_client() -> storage.Client:
    """Initializes and returns a thread-safe GCS client."""
    global gcs_client
    if gcs_client is None:
        try:
            logger.info("Initializing Google Cloud Storage client.")
            gcs_client = storage.Client()
        except Exception as e:
            logger.error(f"FATAL: Failed to initialize GCS client: {e}", exc_info=True)
            raise
    return gcs_client

# --- Core Application Logic ---

def initialize_embedding_model():
    """Loads the sentence-transformer model into memory."""
    global embedding_tokenizer, embedding_model
    try:
        logger.info(f"Loading embedding model '{EMBEDDING_MODEL}' onto device '{device}'")
        embedding_tokenizer = AutoTokenizer.from_pretrained(EMBEDDING_MODEL)
        embedding_model = AutoModel.from_pretrained(EMBEDDING_MODEL).to(device).eval()
        logger.info("Embedding model loaded successfully.")
    except Exception as e:
        logger.error(f"Failed to load embedding model: {e}", exc_info=True)
        raise RuntimeError(f"Could not initialize embedding model: {e}")

def average_pool(last_hidden_states, attention_mask):
    """Pools the output of the transformer model to get a single sentence embedding."""
    last_hidden = last_hidden_states.masked_fill(~attention_mask[..., None].bool(), 0.0)
    return last_hidden.sum(dim=1) / attention_mask.sum(dim=1)[..., None]

async def get_embeddings(texts: List[str]) -> np.ndarray:
    """Generates embeddings for a list of texts without blocking the server."""
    if embedding_model is None or embedding_tokenizer is None:
        raise HTTPException(status_code=500, detail="Embedding model not initialized.")
    if not texts:
        return np.array([])

    prefixed_texts = [f"passage: {text.strip()}" for text in texts if text and text.strip()]
    if not prefixed_texts:
        return np.array([])

    batch_dict = embedding_tokenizer(prefixed_texts, max_length=512, padding=True, truncation=True, return_tensors='pt').to(device)

    # This helper function contains the synchronous, blocking code.
    def compute_embeddings_sync(model, batch):
        with torch.no_grad():
            return model(**batch)

    # We run the blocking function in a separate thread, freeing the main loop.
    outputs = await asyncio.to_thread(compute_embeddings_sync, embedding_model, batch_dict)

    # The rest of the function is fast and can run in the main thread
    embeddings = average_pool(outputs.last_hidden_state, batch_dict['attention_mask'])
    embeddings = torch.nn.functional.normalize(embeddings, p=2, dim=1)
    
    return embeddings.cpu().numpy().astype(np.float32)

def initialize_index():
    """Creates a new, empty FAISS index."""
    global faiss_index, chunks, metadata
    logger.info("Initializing new, empty FAISS index.")
    faiss_index = faiss.IndexFlatIP(1024)
    faiss_index = faiss.IndexIDMap(faiss_index)
    chunks = []
    metadata = []

def save_vector_db_to_gcs():
    """Saves the current index and metadata to local files, then uploads to GCS."""
    logger.info("Attempting to save vector database to GCS...")

    if faiss_index is None or faiss_index.ntotal == 0:
        logger.warning("Skipping save to GCS because index is empty.")
        return

    try:
        logger.info(f"Saving FAISS index with {faiss_index.ntotal} vectors locally.")
        faiss.write_index(faiss_index, str(FAISS_INDEX_PATH))

        logger.info(f"Saving metadata for {len(metadata)} chunks locally as JSONL.")
        with jsonlines.open(METADATA_PATH, mode='w') as writer:
            for i, meta_item in enumerate(metadata):
                # Ensure the 'chunk' key is present in the metadata for persistence
                # and then retrieve it back when loading.
                meta_to_write = meta_item.copy()
                meta_to_write["chunk"] = chunks[i]
                writer.write(meta_to_write)

        client = get_gcs_client()
        bucket = client.bucket(GCS_BUCKET_NAME)

        index_blob = bucket.blob(GCS_INDEX_BLOB_NAME)
        index_blob.upload_from_filename(str(FAISS_INDEX_PATH))
        logger.info(f"Uploaded index to gs://{GCS_BUCKET_NAME}/{GCS_INDEX_BLOB_NAME}")

        metadata_blob = bucket.blob(GCS_METADATA_BLOB_NAME)
        metadata_blob.upload_from_filename(str(METADATA_PATH))
        logger.info(f"Uploaded metadata to gs://{GCS_BUCKET_NAME}/{GCS_METADATA_BLOB_NAME}")
    except Exception as e:
        logger.error(f"Failed to save vector DB to GCS: {e}", exc_info=True)

def load_vector_db_from_gcs():
    """Downloads index/metadata from GCS and loads them into memory."""
    global faiss_index, chunks, metadata
    logger.info("Attempting to load vector database from GCS...")

    try:
        client = get_gcs_client()
        bucket = client.bucket(GCS_BUCKET_NAME)

        index_blob = bucket.blob(GCS_INDEX_BLOB_NAME)
        metadata_blob = bucket.blob(GCS_METADATA_BLOB_NAME)

        if not index_blob.exists() or not metadata_blob.exists():
            logger.warning("Vector DB not found in GCS. A new one will be created on first upload.")
            return False

        logger.info(f"Downloading index from gs://{GCS_BUCKET_NAME}/{GCS_INDEX_BLOB_NAME}")
        index_blob.download_to_filename(str(FAISS_INDEX_PATH))

        logger.info(f"Downloading metadata from gs://{GCS_BUCKET_NAME}/{GCS_METADATA_BLOB_NAME}")
        metadata_blob.download_to_filename(str(METADATA_PATH))

        faiss_index = faiss.read_index(str(FAISS_INDEX_PATH))
        chunks = []
        metadata = []
        with jsonlines.open(METADATA_PATH, mode='r') as reader:
            for meta_item in reader:
                chunk = meta_item.pop("chunk", "")
                chunks.append(chunk)
                metadata.append(meta_item)

        if faiss_index is not None and faiss_index.ntotal == len(chunks) and len(chunks) == len(metadata):
            logger.info(f"Successfully loaded FAISS index with {faiss_index.ntotal} vectors from GCS.")
            return True
        else:
            logger.error("GCS vector DB is inconsistent. Index and metadata do not match. Starting fresh.")
            return False

    except Exception as e:
        logger.error(f"Failed to load vector DB from GCS: {e}", exc_info=True)
        return False

def save_settings_to_gcs():
    """Saves the current app_settings dictionary to a JSON file in GCS."""
    global app_settings
    try:
        client = get_gcs_client()
        bucket = client.bucket(GCS_BUCKET_NAME)
        blob = bucket.blob(GCS_SETTINGS_BLOB_NAME)
        
        # Ensure the settings object is serializable
        settings_json = json.dumps(app_settings, indent=2)
        
        blob.upload_from_string(
            settings_json,
            content_type="application/json"
        )
        logger.info(f"Successfully saved settings to gs://{GCS_BUCKET_NAME}/{GCS_SETTINGS_BLOB_NAME}")
    except Exception as e:
        logger.error(f"Failed to save settings to GCS: {e}", exc_info=True)

def load_settings_from_gcs():
    """Loads settings from GCS and updates the global app_settings variable."""
    global app_settings
    try:
        client = get_gcs_client()
        bucket = client.bucket(GCS_BUCKET_NAME)
        blob = bucket.blob(GCS_SETTINGS_BLOB_NAME)
        
        if blob.exists():
            logger.info(f"Loading settings from gs://{GCS_BUCKET_NAME}/{GCS_SETTINGS_BLOB_NAME}")
            settings_data = blob.download_as_text()
            loaded_settings = json.loads(settings_data)
            app_settings.update(loaded_settings)
            logger.info(f"Settings loaded successfully: {app_settings.get('botName')}")
        else:
            logger.warning("Settings file not found in GCS. Using default settings and creating file.")
            # Save the default settings to create the file for the first time
            save_settings_to_gcs()
            
    except Exception as e:
        logger.error(f"Failed to load settings from GCS. Using default settings. Error: {e}", exc_info=True)

# --- Lifespan Management ---
@asynccontextmanager
async def lifespan(app: FastAPI):
    """Manages application startup and shutdown events."""
    logger.info("API Lifespan: Startup sequence initiated.")
    try:
        load_settings_from_gcs()
        initialize_embedding_model()
        if not load_vector_db_from_gcs():
            initialize_index()
        logger.info("API Lifespan: Startup complete. Ready to serve requests.")
        yield
    except Exception as e:
        logger.critical(f"API Lifespan: FATAL STARTUP ERROR: {e}", exc_info=True)
        raise
    finally:
        logger.info("API Lifespan: Shutdown sequence initiated.")
        save_vector_db_to_gcs()
        logger.info("API Lifespan: Shutdown complete.")

app = FastAPI(
    title="SpikedAI Backend API",
    version="4.1.2",
    description="FastAPI based backend for SPIKED AI. Built and maintained by Umar Yaksambi.",
    lifespan=lifespan
)

app.add_middleware(
    CORSMiddleware,
    allow_origins=["*"],
    allow_credentials=True,
    allow_methods=["*"],
    allow_headers=["*"],
)

def extract_questions_from_pdf_gcs(pdf_path: Path) -> list:
    """Downloads a PDF from GCS and extracts lines/questions from it."""
    questions = []
    try:
        client = get_gcs_client()
        bucket = client.bucket(GCS_BUCKET_NAME)
        blob = bucket.blob(str(pdf_path))
        if not blob.exists():
            logger.warning(f"Question bank PDF not found in GCS at {pdf_path}. No golden questions loaded.")
            return []
        with tempfile.NamedTemporaryFile(suffix=".pdf", delete=True) as tmp:
            blob.download_to_filename(tmp.name)
            doc = fitz.open(tmp.name)
            for page in doc:
                text = page.get_textpage().extractText()
                lines = [line.strip() for line in text.splitlines() if line.strip()]
                questions.extend(lines)
            doc.close()
    except Exception as e:
        logger.error(f"Failed to extract questions from GCS PDF: {e}")
    return questions

# --- At startup ---
QUESTION_BANK_GCS_PATH = "Question_Bank/question_bank.pdf"
QUESTION_BANK = extract_questions_from_pdf_gcs(Path(QUESTION_BANK_GCS_PATH))
logger.info(f"Loaded {len(QUESTION_BANK)} questions from gs://{GCS_BUCKET_NAME}/{QUESTION_BANK_GCS_PATH}")

# --- Pydantic Models ---
class AskRequest(BaseModel):
    question: str = Field(..., min_length=1, max_length=1000)

class UploadResponse(BaseModel):
    message: str
    filename: str
    chunks_added: int
    total_chunks: int

class AskResponse(BaseModel):
    answer: str
    question: str
    sources: List[dict]
    client_followups: List[str] = Field(default_factory=list)
    sales_followups: List[str] = Field(default_factory=list)

class CrawlRequest(BaseModel):
    url: HttpUrl

class CrawlResponse(BaseModel):
    url: HttpUrl
    title: str | None = None
    author: str | None = None
    date: str | None = None
    hostname: str | None = None
    main_text: str
    comments: str | None = None
    raw_text: str | None = None
    source: str | None = None

class WebsiteURLResponse(BaseModel):
    source_url: HttpUrl
    filename: str

class TranscriptEntry(BaseModel):
    speaker: Optional[str] = None
    text: str

class GetQuestionsRequest(BaseModel):
    transcript: List[TranscriptEntry]
    auto_questions_enabled: bool = Field(..., alias="auto_questions_enabled")
    meeting_active: bool = Field(..., alias="meeting_active")

class GetQuestionsResponse(BaseModel):
    questions: List[str]

class SettingsModel(BaseModel):
    botName: str
    selectedPersona: str
    selectedAnswerStyles: List[str]
    customPrompt: str

# --- Text Extraction and Processing ---

def chunk_text(text: str) -> List[str]:
    """Chunks text into smaller pieces using a tokenizer."""
    if not text: return []
    tokens = tokenizer.encode(text)
    chunks_by_tokens = [tokens[i:i + CHUNK_SIZE] for i in range(0, len(tokens), CHUNK_SIZE - OVERLAP_SIZE)]
    return [tokenizer.decode(chunk) for chunk in chunks_by_tokens]

def extract_text_from_file(file_path: Path) -> str:
    """Extracts text from various file types, using OCR as a fallback for PDFs."""
    ext = file_path.suffix.lower()
    text = ""
    try:
        if ext == ".pdf":
            doc = fitz.Document(str(file_path))
            for page in doc:
                page_text = page.get_textpage().extractText()
                page_num = page.number if page.number is not None else 0
                if page_text and page_text.strip():
                    text += f"\n--- Page {page_num + 1} ---\n{page_text}"
                else: # Fallback to OCR for image-based pages
                    images = convert_from_path(str(file_path), first_page=page_num + 1, last_page=page_num + 1)
                    if images:
                        ocr_text = pytesseract.image_to_string(images[0])
                        if ocr_text and ocr_text.strip():
                            text += f"\n--- Page {page_num + 1} [OCR] ---\n{ocr_text}"
        elif ext == ".docx":
            doc = docx.Document(str(file_path))
            text = "\n".join([para.text for para in doc.paragraphs if para.text.strip()])
        elif ext in [".xlsx", ".xls"]:
            wb = openpyxl.load_workbook(file_path, data_only=True)
            for sheet in wb.worksheets:
                for row in sheet.iter_rows():
                    row_text = [str(cell.value) for cell in row if cell.value is not None]
                    if any(cell.strip() for cell in row_text): text += "\t".join(row_text) + "\n"
        elif ext in [".pptx", ".ppt"]:
            prs = Presentation(str(file_path))
            for slide in prs.slides:
                slide_text = [
                    shape.text_frame.text # type: ignore
                    for shape in slide.shapes
                    if hasattr(shape, "text_frame") and shape.text_frame is not None and hasattr(shape.text_frame, "text") and shape.text_frame.text and shape.text_frame.text.strip() # type: ignore
                ]
                if slide_text: text += "\n".join(slide_text) + "\n"
        else:
            raise HTTPException(status_code=400, detail=f"Unsupported file format: {ext}")

        if not text.strip():
            raise HTTPException(status_code=400, detail=f"No extractable text found in {file_path.name}")
        return text.strip()

    except Exception as e:
        logger.error(f"Failed to extract text from {file_path.name}: {e}", exc_info=True)
        if isinstance(e, HTTPException): raise
        raise HTTPException(status_code=500, detail=f"Error processing file {file_path.name}")

async def rebuild_index_and_cleanup_task(sanitized_filename: str):
    """
    This function runs in the background. It contains the slow logic to rebuild
    the FAISS index and clean up cloud storage after a document is deleted.
    """
    # This line is the critical fix
    global faiss_index, chunks, metadata
    
    logger.info(f"Background task started: Deleting '{sanitized_filename}' and rebuilding index.")

    # ... (rest of the function is identical to the one I sent before)
    ids_to_remove = {meta['id'] for meta in metadata if meta.get('filename') == sanitized_filename}
    if not ids_to_remove:
        logger.warning(f"Background task for '{sanitized_filename}' found no metadata to delete. Aborting.")
        return

    remaining_chunks = [chunk for i, chunk in enumerate(chunks) if metadata[i]['id'] not in ids_to_remove]
    remaining_metadata = [meta for meta in metadata if meta['id'] not in ids_to_remove]

    if not remaining_chunks:
        logger.info(f"No chunks remaining after deleting '{sanitized_filename}'. Resetting index.")
        initialize_index()
    else:
        logger.info(f"Re-embedding {len(remaining_chunks)} chunks and rebuilding FAISS index.")
        remaining_embeddings = await get_embeddings(remaining_chunks)
        if remaining_embeddings.size == 0:
            logger.error(f"Failed to generate embeddings for remaining chunks. Aborting background delete for {sanitized_filename}.")
            return

        new_faiss_index = faiss.IndexFlatIP(1024)
        new_faiss_index = faiss.IndexIDMap(new_faiss_index)
        new_ids = np.arange(len(remaining_chunks)).astype('int64')
        new_faiss_index.add_with_ids(remaining_embeddings, new_ids)

        for i, meta_item in enumerate(remaining_metadata):
            meta_item['id'] = i
            
        # Now this assignment correctly modifies the global variables
        faiss_index = new_faiss_index
        chunks = remaining_chunks
        metadata = remaining_metadata

    save_vector_db_to_gcs()

    try:
        client = get_gcs_client()
        bucket = client.bucket(GCS_BUCKET_NAME)
        blob_path = f"{GCS_DOCS_PREFIX}{sanitized_filename}"
        blob = bucket.blob(blob_path)
        if blob.exists():
            blob.delete()
    except Exception as e:
        logger.error(f"Background task failed to delete document from GCS: {e}")

    logger.info(f"Background task finished for '{sanitized_filename}'.")

async def rebuild_after_website_delete_task(website_identifier: str):
    """
    This background task finds all data related to a website, removes it,
    and then safely rebuilds the entire index.
    """
    global faiss_index, chunks, metadata
    logger.info(f"Background task started: Deleting data for '{website_identifier}' and rebuilding index.")

    # 1. Find all metadata and filenames related to the website identifier
    website_metadata_entries = [
        meta for meta in metadata
        if meta.get('source_url') == website_identifier or meta.get('filename') == website_identifier
    ]
    if not website_metadata_entries:
        logger.warning(f"Background task for '{website_identifier}' found no data to delete. Aborting.")
        return

    ids_to_remove = {meta['id'] for meta in website_metadata_entries}
    filenames_to_delete_from_gcs = {meta['filename'] for meta in website_metadata_entries if 'filename' in meta}

    # 2. Create new lists containing only the data we want to KEEP
    remaining_chunks = [chunk for i, chunk in enumerate(chunks) if metadata[i]['id'] not in ids_to_remove]
    remaining_metadata = [meta for meta in metadata if meta['id'] not in ids_to_remove]

    # 3. Perform a full re-index with the remaining data
    if not remaining_chunks:
        logger.info(f"No chunks remaining after deleting '{website_identifier}'. Resetting index.")
        initialize_index()
    else:
        logger.info(f"Re-embedding {len(remaining_chunks)} chunks and rebuilding FAISS index.")
        remaining_embeddings = await get_embeddings(remaining_chunks)
        if remaining_embeddings.size == 0:
            logger.error(f"Failed to generate embeddings for remaining chunks. Aborting background delete for {website_identifier}.")
            return

        new_faiss_index = faiss.IndexFlatIP(1024)
        new_faiss_index = faiss.IndexIDMap(new_faiss_index)
        new_ids = np.arange(len(remaining_chunks)).astype('int64')
        new_faiss_index.add_with_ids(remaining_embeddings, new_ids)

        for i, meta_item in enumerate(remaining_metadata):
            meta_item['id'] = i

        faiss_index = new_faiss_index
        chunks = remaining_chunks
        metadata = remaining_metadata

    # 4. Save the new, correct state to GCS
    save_vector_db_to_gcs()

    # 5. Delete the original source files from GCS
    try:
        client = get_gcs_client()
        bucket = client.bucket(GCS_BUCKET_NAME)
        for fname in filenames_to_delete_from_gcs:
            blob_path = f"{GCS_DOCS_PREFIX}{fname}"
            blob = bucket.blob(blob_path)
            if blob.exists():
                blob.delete()
                logger.info(f"Deleted original crawled document from GCS: {blob_path}")
    except Exception as e:
        logger.error(f"Background task failed to delete one or more source documents from GCS: {e}")

    logger.info(f"Background task for '{website_identifier}' finished.")

async def query_groq(messages: List[dict]) -> str:
    try:
        async with httpx.AsyncClient(timeout=120.0) as client:
            response = await client.post(
                GROQ_API_URL,
                headers={"Authorization": f"Bearer {GROQ_API_KEY}"},
                json={"model": GROQ_MODEL, "messages": messages, "temperature": 0.2, "max_tokens": 800}
            )
            logger.info(f"Groq API response status: {response.status_code}")
            logger.info(f"Groq API response body: {response.text}")
            response.raise_for_status()
            return response.json()["choices"][0]["message"]["content"].strip()
    except Exception as e:
        logger.error(f"Groq query failed: {e}")
        raise HTTPException(status_code=502, detail="Failed to get response from language model.")

def find_golden_questions(followups: list, question_bank: list, threshold: float = 0.85) -> tuple[list, list]:
    """Returns (golden_questions, other_questions) based on similarity to the question bank."""
    golden, others = [], []
    for fq in followups:
        matches = difflib.get_close_matches(fq, question_bank, n=1, cutoff=threshold)
        if matches or fq in question_bank:
            golden.append(fq)
        else:
            others.append(fq)
    return golden, others

def sanitize_filename(filename: str) -> str:
    filename = os.path.basename(filename)
    filename = re.sub(r'[^a-zA-Z0-9_\-.]', '_', filename)
    return filename

# --- Vexa API Integration ---
def get_vexa_headers():
    if not VEXA_API_KEY:
        raise HTTPException(status_code=500, detail="VEXA_API_KEY not set in environment.")
    return {
        "X-API-Key": VEXA_API_KEY,
        "Content-Type": "application/json"
    }

async def vexa_request(method: str, endpoint: str, params=None, json=None):
    url = f"{VEXA_API_BASE_URL}{endpoint}"
    headers = get_vexa_headers()
    async with httpx.AsyncClient(timeout=60.0) as client:
        resp = await client.request(method, url, headers=headers, params=params, json=json)
        try:
            resp.raise_for_status()
        except Exception as e:
            logger.error(f"Vexa API error: {resp.text}")
            raise HTTPException(status_code=resp.status_code, detail=resp.text)
        return resp.json()

# --- Vexa API Endpoints ---

@app.get("/settings", response_model=SettingsModel, summary="Get current AI settings", tags=["Settings"])
async def get_settings():
    """
    Retrieves the current AI personality and style settings that are being used
    by the /ask endpoint.
    """
    return app_settings

@app.post("/settings", response_model=SettingsModel, summary="Update AI settings", tags=["Settings"])
async def update_settings(settings: SettingsModel, background_tasks: BackgroundTasks):
    """
    Updates and saves the AI personality and style settings. The changes are
    persisted to Google Cloud Storage and will be reflected in subsequent
    calls to the /ask endpoint.
    """
    global app_settings
    app_settings = settings.dict()
    background_tasks.add_task(save_settings_to_gcs)
    logger.info(f"Settings updated to: {app_settings}")
    return app_settings

@app.post("/vexa/bots", summary="Request a Vexa bot for a meeting", tags=["Vexa"])
async def vexa_request_bot(payload: dict):
    """
    Proxy to Vexa: Request a transcription bot for a meeting.
    Payload should include: platform, native_meeting_id, language (optional), bot_name (optional)
    """
    return await vexa_request("POST", "/bots", json=payload)

@app.get("/vexa/transcripts/{platform}/{meeting_id}", summary="Get real-time meeting transcript from Vexa", tags=["Vexa"])
async def vexa_get_transcript(platform: str, meeting_id: str):
    endpoint = f"/transcripts/{platform}/{meeting_id}"
    return await vexa_request("GET", endpoint)

@app.get("/vexa/bots/status", summary="Get status of running Vexa bots", tags=["Vexa"])
async def vexa_bots_status():
    return await vexa_request("GET", "/bots/status")

@app.put("/vexa/bots/{platform}/{meeting_id}/config", summary="Update Vexa bot configuration", tags=["Vexa"])
async def vexa_update_bot_config(platform: str, meeting_id: str, payload: dict):
    endpoint = f"/bots/{platform}/{meeting_id}/config"
    return await vexa_request("PUT", endpoint, json=payload)

@app.delete("/vexa/bots/{platform}/{meeting_id}", summary="Stop a Vexa bot", tags=["Vexa"])
async def vexa_stop_bot(platform: str, meeting_id: str):
    endpoint = f"/bots/{platform}/{meeting_id}"
    return await vexa_request("DELETE", endpoint)

@app.get("/vexa/meetings", summary="List your Vexa meetings", tags=["Vexa"])
async def vexa_list_meetings():
    return await vexa_request("GET", "/meetings")

@app.patch("/vexa/meetings/{platform}/{meeting_id}", summary="Update Vexa meeting data", tags=["Vexa"])
async def vexa_update_meeting(platform: str, meeting_id: str, payload: dict):
    endpoint = f"/meetings/{platform}/{meeting_id}"
    return await vexa_request("PATCH", endpoint, json=payload)

@app.delete("/vexa/meetings/{platform}/{meeting_id}", summary="Delete Vexa meeting and transcripts", tags=["Vexa"])
async def vexa_delete_meeting(platform: str, meeting_id: str):
    endpoint = f"/meetings/{platform}/{meeting_id}"
    return await vexa_request("DELETE", endpoint)

@app.put("/vexa/user/webhook", summary="Set Vexa user webhook URL", tags=["Vexa"])
async def vexa_set_webhook(payload: dict):
    endpoint = "/user/webhook"
    return await vexa_request("PUT", endpoint, json=payload)

# --- API Endpoints ---
@app.get("/", summary="Health Check", tags=["System"])
def read_root():
    """A simple endpoint to confirm the API is running."""
    return {"status": "Sales Assistant API is online and running."}

@app.post("/refresh", summary="Synchronize the index with GCS documents", tags=["Document Management"])
async def refresh_index():
    """
    Checks for inconsistencies and performs a full re-index if needed.
    """
    global faiss_index, chunks, metadata
    logger.info("Starting a comprehensive refresh of the search index.")

    client = get_gcs_client()
    bucket = client.bucket(GCS_BUCKET_NAME)
    
    gcs_blobs = list(bucket.list_blobs(prefix=GCS_DOCS_PREFIX))
    gcs_filenames = {Path(blob.name).name for blob in gcs_blobs if not blob.name.endswith("/")}
    indexed_filenames = {meta["filename"] for meta in metadata if "filename" in meta}
    
    new_files_to_add = gcs_filenames - indexed_filenames
    stale_files_to_remove = indexed_filenames - gcs_filenames
    is_count_mismatched = (faiss_index.ntotal != len(chunks) or len(chunks) != len(metadata))
    
    full_reindex_needed = bool(stale_files_to_remove) or is_count_mismatched

    if full_reindex_needed:
        # --- Corrected Logging ---
        logger.warning("Inconsistency detected! Performing a full re-index for the following reasons:")
        if stale_files_to_remove:
            logger.warning(f"- Stale files found in index but not in GCS: {stale_files_to_remove}")
        if is_count_mismatched:
            logger.warning(f"- Mismatched counts: FAISS={faiss_index.ntotal}, Chunks={len(chunks)}, Metadata={len(metadata)}")
        # --- End Corrected Logging ---

        # ... (The rest of the full re-index logic is the same)
        initialize_index()
        all_chunks = []
        all_metadata = []
        
        for blob in gcs_blobs:
            filename = Path(blob.name).name
            if not filename: continue

            logger.info(f"[Full Re-index] Processing: {filename}")
            with tempfile.TemporaryDirectory() as temp_dir:
                temp_path = Path(temp_dir) / filename
                try:
                    blob.download_to_filename(str(temp_path))
                    text = extract_text_from_file(temp_path)
                    text_chunks = chunk_text(text)
                    
                    if text_chunks:
                        for chunk in text_chunks:
                            all_chunks.append(chunk)
                            all_metadata.append({"filename": filename})
                except Exception as e:
                    logger.error(f"[Full Re-index] Failed to process {filename}: {e}", exc_info=True)

        if not all_chunks:
            logger.info("Full re-index complete. No documents found to index.")
            save_vector_db_to_gcs()
            return {"message": "Full re-index complete. No documents found in GCS."}

        logger.info(f"Generating embeddings for {len(all_chunks)} total chunks... (This may take several minutes on a CPU)")
        all_embeddings = await get_embeddings(all_chunks)
        new_ids = np.arange(len(all_chunks)).astype("int64")
        faiss_index.add_with_ids(all_embeddings, new_ids)
        
        for i, meta in enumerate(all_metadata):
            meta['id'] = i
        
        chunks = all_chunks
        metadata = all_metadata
        
        save_vector_db_to_gcs()
        logger.info(f"Full re-index complete. Total chunks: {len(chunks)}")
        return {"message": f"Full re-index complete. Processed {len(gcs_filenames)} documents. Total chunks: {len(chunks)}."}

    elif new_files_to_add:
        # (This part of the function remains the same as before)
        logger.info(f"No inconsistencies found. Adding {len(new_files_to_add)} new document(s).")
        # ... (rest of the incremental add logic)
        new_count = 0
        blobs_to_add = [b for b in gcs_blobs if Path(b.name).name in new_files_to_add]
        for blob in blobs_to_add:
            with tempfile.TemporaryDirectory() as temp_dir:
                temp_path = Path(temp_dir) / Path(blob.name).name
                blob.download_to_filename(str(temp_path))
                try:
                    text = extract_text_from_file(temp_path)
                    text_chunks = chunk_text(text)
                    if not text_chunks:
                        continue
                    embeddings = await get_embeddings(text_chunks)
                    if embeddings.size == 0:
                        continue
                    start_index = len(chunks)
                    ids = np.arange(start_index, start_index + len(text_chunks)).astype("int64")
                    faiss_index.add_with_ids(embeddings, ids)
                    for i, chunk in enumerate(text_chunks):
                        chunks.append(chunk)
                        metadata.append({"id": start_index + i, "filename": Path(blob.name).name})
                    new_count += 1
                except Exception as e:
                    logger.error(f"Error processing new file {blob.name}: {e}", exc_info=True)
        if new_count > 0:
            save_vector_db_to_gcs()
        return {"message": f"Refresh complete. Indexed {new_count} new document(s). Total chunks: {len(chunks)}"}

    else:
        logger.info("Refresh complete. No changes detected. Index is up to date.")
        return {"message": "All GCS documents are already indexed and consistent."}

@app.post("/upload", response_model=UploadResponse, summary="Upload and Process a Document", tags=["Search AI"])
async def upload_file(file: UploadFile = File(...)):
    try:
        filename = file.filename if file.filename is not None else ""
        sanitized_filename = sanitize_filename(filename)
        if not sanitized_filename or Path(sanitized_filename).suffix.lower() not in ALLOWED_EXTENSIONS:
            raise HTTPException(status_code=400, detail="Unsupported or unsafe file type.")

        # Check for duplicates
        if any(meta['filename'] == sanitized_filename for meta in metadata):
            raise HTTPException(status_code=409, detail=f"Document '{sanitized_filename}' already exists. Please delete it first if you want to re-upload.")

        with tempfile.TemporaryDirectory() as temp_dir:
            file_path = Path(temp_dir) / sanitized_filename
            with open(file_path, "wb") as f:
                f.write(await file.read())
            logger.info(f"Temporarily saved uploaded file to {file_path}")

            text = extract_text_from_file(file_path)
            text_chunks = chunk_text(text)
            if not text_chunks:
                raise HTTPException(status_code=400, detail="Failed to create text chunks from document.")

            new_embeddings = await get_embeddings(text_chunks)
            if new_embeddings.size == 0:
                raise HTTPException(status_code=500, detail="Failed to generate embeddings.")

            start_index = len(chunks)
            ids_to_add = np.arange(start_index, start_index + len(text_chunks)).astype('int64')
            faiss_index.add_with_ids(new_embeddings, ids_to_add) # type: ignore

            for i, chunk in enumerate(text_chunks):
                chunk_id = start_index + i
                chunks.append(chunk)
                metadata.append({"id": chunk_id, "filename": sanitized_filename})

            save_vector_db_to_gcs()

            client = get_gcs_client()
            bucket = client.bucket(GCS_BUCKET_NAME)
            blob = bucket.blob(f"{GCS_DOCS_PREFIX}{sanitized_filename}")
            blob.upload_from_filename(str(file_path))
            logger.info(f"Uploaded original document to gs://{GCS_BUCKET_NAME}/{blob.name}")

            return UploadResponse(
                message="File uploaded and processed successfully",
                filename=sanitized_filename,
                chunks_added=len(text_chunks),
                total_chunks=len(chunks)
            )
    except HTTPException as e:
        logger.error(f"HTTP exception during upload: {e.detail}")
        raise
    except Exception as e:
        logger.error(f"Unexpected error during upload: {e}", exc_info=True)
        raise HTTPException(status_code=500, detail="An internal server error occurred.")

@app.post("/ask", response_model=AskResponse, summary="Ask a Question and get relevant information and follow-up questions")
async def ask_question(request: AskRequest):
    """
    Searches for relevant context from uploaded documents, uses an LLM to answer the question,
    and generates two types of follow-up questions: for the sales person and for the client.
    The AI's personality and response style are based on the saved settings.
    """

    question = request.question.strip()
    if not question:
        raise HTTPException(status_code=400, detail="Question cannot be empty.")
    if faiss_index is None or faiss_index.ntotal == 0:
        raise HTTPException(status_code=404, detail="No documents uploaded or indexed yet.")

    query_embedding = await get_embeddings([f"query: {question}"])
    if query_embedding.size == 0:
        raise HTTPException(status_code=500, detail="Could not generate embedding for the question.")

    k = min(TOP_K, faiss_index.ntotal)
    distances, labels = faiss_index.search(query_embedding, k)  # type: ignore

    relevant_chunks, sources = [], set()
    for i, chunk_id in enumerate(labels[0]):
        if chunk_id != -1 and distances[0][i] > 0.3:
            relevant_chunks.append(chunks[chunk_id])
            if 'filename' in metadata[chunk_id]:
                sources.add(metadata[chunk_id]["filename"])

    clean_answer = "I couldn't find relevant information in the uploaded documents to answer your question."
    sales_followups_list = []
    client_followups_list = []

    if relevant_chunks:
        context = "\n\n".join(relevant_chunks)
        
        # --- DYNAMIC PROMPT INJECTION ---
        custom_prompt_section = app_settings.get("customPrompt", "")
        
        base_system_prompt = """
Your responsibilities:
1. Provide detailed answers to questions using the provided context.
2. Structure the response clearly by breaking it into multiple paragraphs as needed. Use appropriate formatting to highlight key points, such as bold text for emphasis, bullet points for lists, and headings for sectioning.
3. If the context is insufficient, state so clearly.
4. If the question includes general knowledge, use it to enhance your answer.
5. After your answer, generate two sets of follow-up questions (each as a bullet list):
   - Sales Follow-up Questions: Questions the sales person can ask the client, answerable using the context but not already covered in the answer.
        Start with phrases like:
            Would you like to know...,
            Should I walk you through...,
            Do you want to explore...,
            Can I show you...,
            Are you interested in...
   - Client Follow-up Questions: Questions the client might ask the sales person, answerable using the context but not already covered in the answer.
6. Never hallucinate. If information is missing, say so.
7. If no context is provided, respond accordingly.

Structure your response as:
Answer:
<your answer>

Sales Follow-up Questions:
- <question 1>
- <question 2>
...

Client Follow-up Questions:
- <question 1>
- <question 2>
...
"""
        
        final_system_prompt = f"{custom_prompt_section}\n\n{base_system_prompt}"
        
        messages = [
            {"role": "system", "content": final_system_prompt},
            {"role": "user", "content": f"Context:\n{context}\n\nQuestion: {question}"}
        ]

        raw_response = await query_groq(messages)

        # Clean answer extraction
        answer_match = re.split(r"(?i)\n\s*sales follow[- ]?up questions[:\-]?\s*\n", raw_response)
        clean_answer = answer_match[0].strip() if answer_match else raw_response.strip()

        # Extract sales follow-ups
        sales_followups = []
        client_followups = []

        sales_match = re.search(
            r"(?i)sales follow[- ]?up questions[:\-]?\s*(.*?)(?:(?:client follow[- ]?up questions[:\-]?)|$)",
            raw_response,
            re.DOTALL
        )
        if sales_match:
            sales_lines = sales_match.group(1).strip().splitlines()
            sales_followups = [line.strip(" -•1234567890.").strip() for line in sales_lines if line.strip()]

        # Extract client follow-ups
        client_match = re.search(
            r"(?i)client follow[- ]?up questions[:\-]?\s*(.*)", raw_response, re.DOTALL
        )
        if client_match:
            client_lines = client_match.group(1).strip().splitlines()
            client_followups = [line.strip(" -•1234567890.").strip() for line in client_lines if line.strip()]

        # Limit to 3 each
        sales_followups = sales_followups[:3]
        client_followups = client_followups[:3]

    return AskResponse(
        answer=clean_answer,
        question=question,
        sources=[
        {
            "filename": fname,
            "url": f"{BASE_URL}/download/{quote(fname)}"        
        }
        for fname in sources
        ],
        client_followups=client_followups,
        sales_followups=sales_followups
    )

@app.post("/getQuestions", response_model=GetQuestionsResponse, summary="Generate relevant questions from a transcript", tags=["Search AI"])
async def get_questions(request: GetQuestionsRequest):
    """
    Analyzes the latest part of a meeting transcript and generates relevant questions
    that could be asked next, based on the conversation's context.
    """
    if not request.meeting_active or not request.auto_questions_enabled or not request.transcript:
        return GetQuestionsResponse(questions=[])

    try:
        # Use the last 10 messages for context, as in the reference
        recent_messages = request.transcript[-10:]
        formatted_transcript = "\n".join(
            [f"{msg.speaker or 'Unknown'}: {msg.text}" for msg in recent_messages]
        )

        # Construct the prompt for the LLM
        prompt_content = (
            "Based on this meeting conversation, generate 3 relevant questions that would be valuable to ask next:"
            f"\n\n{formatted_transcript}\n\n"
            "Format as: 1. [Question]\n2. [Question]\n3. [Question]"
        )

        messages = [{"role": "user", "content": prompt_content}]
        raw_response = await query_groq(messages)

        # Parse the response to extract questions
        lines = raw_response.strip().splitlines()
        questions = []
        for line in lines:
            # Match lines that start with a number and a period (e.g., "1.", "2. ")
            if re.match(r"^\d+\.\s*", line.strip()):
                # Remove the leading number and period to get the clean question text
                question_text = re.sub(r"^\d+\.\s*", "", line.strip()).strip()
                if question_text:
                    questions.append(question_text)

        # Limit to the top 3 questions
        final_questions = questions[:3]

        return GetQuestionsResponse(questions=final_questions)

    except HTTPException:
        # Re-raise HTTP exceptions from query_groq
        raise
    except Exception as e:
        logger.error(f"Error generating questions from transcript: {e}", exc_info=True)
        # Return an empty list or an error response
        raise HTTPException(status_code=500, detail="Failed to generate questions from transcript.")


@app.get("/documents", summary="List all indexed documents and metadata")
def list_documents():
    docs = {}
    for meta in metadata:
        fname = meta.get("filename")
        chunk_id = meta.get("id")
        if fname and chunk_id is not None:
            docs.setdefault(fname, []).append(chunk_id)
    return [{"filename": fname, "chunk_ids": ids, "total_chunks": len(ids)} for fname, ids in docs.items()]

@app.delete("/documents/{filename}", summary="Delete a Document and its Embeddings", tags=["Document Management"])
async def delete_document(filename: str, background_tasks: BackgroundTasks):
    """
    Responds immediately and schedules a background task to delete a document
    and rebuild the search index.
    """
    sanitized_filename = sanitize_filename(filename)

    if not any(meta.get('filename') == sanitized_filename for meta in metadata):
        raise HTTPException(status_code=404, detail=f"Document '{sanitized_filename}' not found.")

    background_tasks.add_task(rebuild_index_and_cleanup_task, sanitized_filename)

    return {
        "message": f"Deletion of '{sanitized_filename}' initiated. Index rebuild will run in the background."
    }

@app.get("/download/{filename}", summary="Download or view a document from GCS", tags=["Document Management"])
async def download_document(filename: str, background_tasks: BackgroundTasks):
    """
    Streams PDF files for in-browser viewing.
    For all other file types, it streams the content as a download attachment.
    """
    sanitized_filename = sanitize_filename(filename)
    blob_path = f"{GCS_DOCS_PREFIX}{sanitized_filename}"

    try:
        client = get_gcs_client()
        bucket = client.bucket(GCS_BUCKET_NAME)
        blob = bucket.blob(blob_path)

        if not blob.exists():
            raise HTTPException(status_code=404, detail=f"File '{sanitized_filename}' not found.")

        # --- Logic to stream PDFs and download other files ---
        if sanitized_filename.lower().endswith('.pdf'):
            return StreamingResponse(
                blob.open("rb"),
                media_type="application/pdf",
                headers={"Content-Disposition": f'inline; filename="{sanitized_filename}"'}
            )
        else:
            # Stream directly from GCS blob for other file types as attachment
            # Set a more appropriate media type if known, otherwise application/octet-stream
            media_type = "application/octet-stream"
            # Attempt to infer a more specific media type based on extension
            if sanitized_filename.lower().endswith('.docx'):
                media_type = "application/vnd.openxmlformats-officedocument.wordprocessingml.document"
            elif sanitized_filename.lower().endswith('.xlsx'):
                media_type = "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet"
            elif sanitized_filename.lower().endswith('.pptx'):
                media_type = "application/vnd.openxmlformats-officedocument.presentationml.presentation"

            return StreamingResponse(
                blob.open("rb"),
                media_type=media_type,
                headers={"Content-Disposition": f'attachment; filename="{sanitized_filename}"'}
            )

    except HTTPException:
        raise
    except Exception as e:
        logger.error(f"Failed to download file {sanitized_filename} from GCS: {e}", exc_info=True)
        raise HTTPException(status_code=500, detail="Failed to download file from cloud storage.")

@app.delete("/reset", summary="Delete all documents and vector DB", tags=["Document Management"])
def delete_all_documents():
    initialize_index()
    save_vector_db_to_gcs()

    client = get_gcs_client()
    bucket = client.bucket(GCS_BUCKET_NAME)

    deleted_blobs = []
    try:
        # List and delete all blobs under the GCS_DOCS_PREFIX
        blobs_to_delete = list(bucket.list_blobs(prefix=GCS_DOCS_PREFIX))
        for blob in blobs_to_delete:
            blob.delete()
            deleted_blobs.append(blob.name)

        logger.info(f"Deleted {len(deleted_blobs)} documents from GCS.")

        # Ensure vector store blobs are also deleted from GCS
        for vs_blob_name in [GCS_INDEX_BLOB_NAME, GCS_METADATA_BLOB_NAME]:
            vs_blob = bucket.blob(vs_blob_name)
            if vs_blob.exists():
                vs_blob.delete()
                deleted_blobs.append(vs_blob_name)
                logger.info(f"Deleted vector store component from GCS: {vs_blob_name}")

    except Exception as e:
        logger.error(f"Error deleting GCS documents during reset: {e}", exc_info=True)
        raise HTTPException(status_code=500, detail="Failed to delete all documents from cloud storage.")

    return {"message": "All documents and vector DB deleted.", "deleted_files": deleted_blobs}

@app.get("/chunks", summary="Get all chunks and their metadata", tags=["Document Management"])
def get_all_chunks():
    """
    Returns all chunks and their metadata as a JSON array.
    Each item contains: id, filename, chunk
    """
    results = []
    # Iterate over the in-memory metadata and chunks lists
    for i, meta_item in enumerate(metadata):
        # Ensure 'id' and 'filename' exist in metadata, and 'chunks' has the corresponding index
        chunk_id = meta_item.get("id")
        filename = meta_item.get("filename")
        if chunk_id is not None and filename is not None and i < len(chunks):
            results.append({
                "id": chunk_id,
                "filename": filename,
                "chunk": chunks[i]
            })
    return results

# --- Helper Functions ---

async def fetch_html(url: str) -> str:
    headers = {
        "User-Agent": (
            "Mozilla/5.0 (Windows NT 10.0; Win64; x64) AppleWebKit/537.36 "
            "(KHTML, like Gecko) Chrome/114.0.0.0 Safari/537.36"
        )
    }
    try:
        async with httpx.AsyncClient(timeout=20.0, follow_redirects=True) as client:
            response = await client.get(url, headers=headers)
            response.raise_for_status()
            return response.text
    except httpx.TimeoutException:
        logger.error(f"Timeout while fetching URL: {url}")
        raise HTTPException(status_code=408, detail="Timeout while fetching the URL.")
    except httpx.RequestError as e:
        logger.error(f"Network error: {e}")
        raise HTTPException(status_code=400, detail=f"Request error: {e}")
    except httpx.HTTPStatusError as e:
        logger.error(f"HTTP error {e.response.status_code} while fetching {url}")
        raise HTTPException(status_code=e.response.status_code, detail="Bad response from server.")

def fallback_metadata(html: str) -> dict:
    """Fallback parser using BeautifulSoup if Trafilatura doesn't provide full metadata"""
    soup = BeautifulSoup(html, "html.parser")
    title = soup.title.string.strip() if soup.title and soup.title.string else None

    # Attempt to extract author from common meta tags or bylines
    author = None
    author_meta = soup.find('meta', attrs={'name': ['author', 'citation_author']})
    if author_meta and author_meta.get('content'):
        author = author_meta['content']
    elif soup.find(class_=re.compile(r'author|byline', re.I)):
        author_tag = soup.find(class_=re.compile(r'author|byline', re.I))
        if author_tag and author_tag.get_text(strip=True):
            author = author_tag.get_text(strip=True)

    # Attempt to extract date from common meta tags or time elements
    date = None
    date_meta = soup.find('meta', attrs={'property': 'article:published_time'}) or \
                soup.find('meta', attrs={'name': 'date'}) or \
                soup.find('meta', attrs={'itemprop': 'datePublished'})
    if date_meta and date_meta.get('content'):
        date = date_meta['content']
    elif soup.find('time') and soup.find('time').get('datetime'):
        date = soup.find('time')['datetime']

    return {
        "title": title,
        "author": author,
        "date": date,
    }

config = use_config()
config.set("DEFAULT", "EXTRACTION_TIMEOUT", "15")

@app.post("/crawl", response_model=CrawlResponse, summary="Extract clean content from a URL and index it", tags=["Website"])
async def crawl_and_extract(request: CrawlRequest = Body(...)):
    url = str(request.url)
    logger.info(f"Received crawl request: {url}")

    html = await fetch_html(url)

    if not html.strip():
        logger.warning("Empty HTML received")
        raise HTTPException(status_code=404, detail="No HTML content found at the URL.")

    extracted = extract(
        html,
        include_comments=True,
        include_tables=True,
        output_format='json',
        config=config
    )

    if not extracted:
        logger.warning(f"Trafilatura extraction failed for: {url}")
        raise HTTPException(status_code=422, detail="Failed to extract main content.")

    try:
        data = json.loads(extracted)
    except json.JSONDecodeError:
        logger.error("JSON decoding error from Trafilatura output")
        raise HTTPException(status_code=500, detail="Failed to parse extracted content.")

    text = data.get("text")
    if not text:
        raise HTTPException(status_code=422, detail="No readable text found from the URL.")

    # Use fallback metadata if missing from trafilatura's output
    fallback_data = fallback_metadata(html)
    data["title"] = data.get("title") or fallback_data.get("title")
    data["author"] = data.get("author") or fallback_data.get("author")
    data["date"] = data.get("date") or fallback_data.get("date")
    # data["raw_text"] = data.get("raw_text") # Trafilatura usually provides this
    # data["source"] = data.get("source") # Trafilatura usually provides this

    # --- Chunking and Embedding ---
    text_chunks = chunk_text(text)
    if not text_chunks:
        raise HTTPException(status_code=400, detail="Failed to create text chunks from crawled content.")

    new_embeddings = await get_embeddings(text_chunks)
    if new_embeddings.size == 0:
        raise HTTPException(status_code=500, detail="Failed to generate embeddings from crawled content.")

    # Generate sanitized filename from URL's hostname and a UUID for uniqueness
    parsed = urlparse(url)
    # Ensure hostname is clean for filename, replace dots with underscores
    clean_hostname = re.sub(r'[^a-zA-Z0-9_\-]', '_', parsed.hostname or "unknown_host")
    sanitized_filename = f"crawled_{clean_hostname}_{uuid4().hex[:8]}.html"

    # Check for duplicate sanitized filename (highly unlikely with UUID, but good practice)
    if any(meta.get('filename') == sanitized_filename for meta in metadata):
        # If by some cosmic chance, retry with a new UUID or add a timestamp
        sanitized_filename = f"crawled_{clean_hostname}_{datetime.now().strftime('%Y%m%d%H%M%S')}_{uuid4().hex[:4]}.html"


    start_index = len(chunks)
    ids_to_add = np.arange(start_index, start_index + len(text_chunks)).astype('int64')
    faiss_index.add_with_ids(new_embeddings, ids_to_add)  # type: ignore

    for i, chunk in enumerate(text_chunks):
        chunk_id = start_index + i
        chunks.append(chunk)
        metadata.append({"id": chunk_id, "filename": sanitized_filename, "source_url": url})

    save_vector_db_to_gcs()

    with tempfile.NamedTemporaryFile(suffix=".html", delete=False) as tmp_html_file:
        tmp_html_file.write(html.encode('utf-8'))
        temp_html_path = Path(tmp_html_file.name)

    client = get_gcs_client()
    bucket = client.bucket(GCS_BUCKET_NAME)
    blob = bucket.blob(f"{GCS_DOCS_PREFIX}{sanitized_filename}")
    blob.upload_from_filename(str(temp_html_path))
    logger.info(f"Uploaded raw crawled HTML to gs://{GCS_BUCKET_NAME}/{blob.name}")
    os.remove(temp_html_path)

    return CrawlResponse(
        url=request.url,
        title=data.get("title"),
        author=data.get("author"),
        date=data.get("date"),
        hostname=parsed.hostname,
        main_text=text,
        comments=data.get("comments"),
        raw_text=data.get("raw_text"),
        source=data.get("source")
    )

@app.get("/website/url/{identifier}", response_model=WebsiteURLResponse, summary="Get the original source URL for a crawled document", tags=["Website"])
def get_website_url(identifier: str):
    """
    Retrieves the original source URL for a crawled document by its filename or the URL itself.
    """
    # Search through the metadata for a match on either filename or source_url
    for meta in metadata:
        is_match = (meta.get('filename') == identifier) or (meta.get('source_url') == identifier)
        if is_match:
            source_url = meta.get('source_url')
            filename = meta.get('filename')
            if source_url and filename:
                 # Found a matching entry, return the source URL and filename
                return WebsiteURLResponse(source_url=source_url, filename=filename)

    # If the loop completes without finding a match, raise a 404 error
    raise HTTPException(status_code=404, detail=f"No source URL found for identifier '{identifier}'.")

@app.delete("/website/{website}", summary="Delete a Website's Embeddings", tags=["Website"])
async def delete_website(website: str, background_tasks: BackgroundTasks):
    """
    Responds immediately and schedules a background task to delete all data
    associated with a crawled website URL or its generated filename.
    """
    # Quick check to see if the website data exists before starting a task
    if not any(meta.get('source_url') == website or meta.get('filename') == website for meta in metadata):
        raise HTTPException(status_code=404, detail=f"Website or document '{website}' not found.")

    # Schedule the background task
    background_tasks.add_task(rebuild_after_website_delete_task, website)

    return {
        "message": f"Deletion for '{website}' initiated. The index will be rebuilt in the background."
    }
    
# --- AI Training --- #

# Configuration
UPLOAD_DIR = Path("uploads")
UPLOAD_DIR.mkdir(exist_ok=True)

DOCUMENT_STORE = Path("document_store")
DOCUMENT_STORE.mkdir(exist_ok=True)

# Note: AUDIO_UPLOAD_DIR no longer needed since we use Web Speech API

# Pydantic Models
class DocumentType(str, Enum):
    proposal = "proposal"
    technical = "technical"
    presentation = "presentation"
    general = "general"

class DifficultyLevel(str, Enum):
    easy = "Easy"
    medium = "Medium"
    hard = "Hard"

class Persona(str, Enum):
    technical_lead = "Technical Lead"
    business_manager = "Business Manager"
    c_suite_executive = "C-Suite Executive"

class DocumentHeading(BaseModel):
    section: str
    title: str
    content_summary: str

class DocumentAnalysis(BaseModel):
    extractedTopics: List[str]
    keyFeatures: List[str]
    documentType: DocumentType
    extractedHeadings: List[DocumentHeading]

class AnalyzeDocumentRequest(BaseModel):
    filename: str

class CompareResponse(BaseModel):
    score: int = Field(ge=0, le=10)
    coverage: str
    key_points_missed: List[str]
    feedback: str
    strengths: List[str]
    improvements: List[str]
    confidence_level: int = Field(ge=0, le=100)
    response_time: int

class CompareAnswersRequest(BaseModel):
    user_answer: str
    ideal_answer: str
    original_question: str
    response_time: int

class GeneratedQuestion(BaseModel):
    question: str
    difficulty_level: DifficultyLevel
    focus_area: str
    document_references: List[str]
    question_id: str

class GenerateQuestionsRequest(BaseModel):
    persona: Persona
    objective: str
    document_filenames: List[str]
    max_questions: int = 15
    difficulty_levels: List[DifficultyLevel] = [DifficultyLevel.easy, DifficultyLevel.medium, DifficultyLevel.hard]

class GenerateQuestionsResponse(BaseModel):
    questions: List[GeneratedQuestion]
    status: str
    document_count: int
    persona: str

class Document(BaseModel):
    filename: str
    chunk_ids: List[str] = []
    total_chunks: int = 0
    upload_date: datetime
    file_size: int
    file_type: str
    extractedTopics: List[str] = []
    keyFeatures: List[str] = []
    documentType: DocumentType = DocumentType.general
    extractedHeadings: List[DocumentHeading] = []

class UploadResponseTrain(BaseModel):
    filename: str
    chunks_added: int
    total_chunks: int

class AskRequestTrain(BaseModel):
    question: str
    include_followups: bool = True
    include_sales_questions: bool = True

class SourceTrain(BaseModel):
    filename: str
    url: Optional[str] = None

class IdealAnswerResponse(BaseModel):
    answer: str
    sources: List[SourceTrain]
    follow_up_questions: List[str]
    sales_followups: List[str]

# In-memory storage
documents_db: Dict[str, Document] = {}
document_chunks: Dict[str, List[str]] = {}
document_content: Dict[str, str] = {}

# Helper Functions
async def call_groq_llm(prompt: str, system_prompt: str = None, max_tokens: int = 1000) -> str:
    """Call Groq LLM API"""
    try:
        messages = []
        if system_prompt:
            messages.append({"role": "system", "content": system_prompt})
        messages.append({"role": "user", "content": prompt})
        
        async with httpx.AsyncClient() as client:
            response = await client.post(
                GROQ_API_URL,  
                headers={
                    "Authorization": f"Bearer {GROQ_API_KEY}",
                    "Content-Type": "application/json"
                },
                json={
                    "model": GROQ_MODEL,  
                    "messages": messages,
                    "max_tokens": max_tokens,
                    "temperature": 0.7
                },
                timeout=60.0
            )
            
        if response.status_code != 200:
            print(f"Groq API error: {response.status_code}, {response.text}")
            raise HTTPException(status_code=500, detail=f"Groq API error: {response.status_code}")
            
        result = response.json()
        return result["choices"][0]["message"]["content"]
        
    except httpx.TimeoutException:
        raise HTTPException(status_code=500, detail="LLM call timeout")
    except Exception as e:
        print(f"LLM call error: {str(e)}")
        raise HTTPException(status_code=500, detail=f"LLM call failed: {str(e)}")

def extract_pdf_text(content: bytes) -> str:
    """Extract text from PDF"""
    try:
        pdf_reader = PyPDF2.PdfReader(BytesIO(content))
        text = ""
        for page in pdf_reader.pages:
            text += page.extract_text()
        return text
    except Exception as e:
        raise HTTPException(status_code=500, detail=f"Error extracting PDF text: {str(e)}")

def extract_docx_text(content: bytes) -> str:
    """Extract text from DOCX"""
    try:
        doc = docx.Document(BytesIO(content))
        text = ""
        for paragraph in doc.paragraphs:
            text += paragraph.text + "\n"
        return text
    except Exception as e:
        raise HTTPException(status_code=500, detail=f"Error extracting DOCX text: {str(e)}")

def extract_text_from_file_train(filename: str, content: bytes) -> str:
    """Extract text content from uploaded file content"""
    if filename.endswith('.pdf'):
        return extract_pdf_text(content)
    elif filename.endswith('.docx'):
        return extract_docx_text(content)
    elif filename.endswith('.txt'):
        return content.decode('utf-8')
    else:
        raise HTTPException(status_code=400, detail="Unsupported file type")

def chunk_text_train(text: str, chunk_size: int = 1000) -> List[str]:
    """Split text into chunks for processing"""
    words = text.split()
    chunks = []
    current_chunk = []
    current_size = 0
    
    for word in words:
        if current_size + len(word) > chunk_size and current_chunk:
            chunks.append(" ".join(current_chunk))
            current_chunk = [word]
            current_size = len(word)
        else:
            current_chunk.append(word)
            current_size += len(word) + 1
    
    if current_chunk:
        chunks.append(" ".join(current_chunk))
    
    return chunks

def parse_json_from_llm_response(response: str) -> dict:
    """Extract JSON from LLM response that might contain markdown or extra text"""
    try:
        # Try to find JSON block in markdown format first
        json_match = re.search(r'```(?:json)?\s*(\{.*?\})\s*```', response, re.DOTALL | re.IGNORECASE)
        if json_match:
            json_str = json_match.group(1)
        else:
            # Try to find raw JSON
            json_match = re.search(r'\{.*\}', response, re.DOTALL)
            if json_match:
                json_str = json_match.group()
            else:
                raise ValueError("No JSON found in response")
        
        return json.loads(json_str)
    except Exception as e:
        print(f"JSON parsing error: {str(e)}")
        print(f"Raw response: {response}")
        raise ValueError(f"JSON parsing failed: {str(e)}")

# API Endpoints

# app.py - NEW CODE

@app.post("/train/upload", response_model=UploadResponseTrain, tags=["AI Training"])
async def train_upload_document(file: UploadFile = File(...)):
    """Uploads a document for AI training sessions."""
    filename = sanitize_filename(file.filename or "unknown.tmp")
    if filename in documents_db:
        raise HTTPException(status_code=409, detail=f"Document '{filename}' already exists.")
    
    content = await file.read()
    if not content:
        raise HTTPException(status_code=400, detail="Uploaded file is empty.")

    try:
        def _extract():
            ext = Path(filename).suffix.lower()
            if ext == '.pdf':
                reader = PyPDF2.PdfReader(BytesIO(content))
                return "".join(page.extract_text() for page in reader.pages if page.extract_text())
            elif ext == '.docx':
                return "\n".join(p.text for p in docx.Document(BytesIO(content)).paragraphs if p.text)
            elif ext == '.txt':
                return content.decode('utf-8', errors='ignore')
            else:
                raise HTTPException(status_code=400, detail="Unsupported file type for training. Use PDF, DOCX, or TXT.")
        
        text_content = await asyncio.to_thread(_extract)
        if not text_content:
            raise HTTPException(status_code=422, detail="No text could be extracted from the document.")

    except HTTPException:
        raise
    except Exception as e:
        logger.error(f"Error extracting text during training upload: {e}", exc_info=True)
        raise HTTPException(status_code=500, detail=f"Error extracting text: {e}")

    chunks = chunk_text_train(text_content)
    doc_info = Document(
        filename=filename,
        chunk_ids=[f"{filename}_chunk_{i}" for i in range(len(chunks))],
        total_chunks=len(chunks),
        upload_date=datetime.now(timezone.utc),
        file_size=len(content),
        file_type=file.content_type or "application/octet-stream"
    )
    documents_db[filename] = doc_info
    document_content[filename] = text_content

    file_path = UPLOAD_DIR / filename
    async with aiofiles.open(file_path, 'wb') as f:
        await f.write(content)
        
    return UploadResponseTrain(
        filename=filename, chunks_added=len(chunks), total_chunks=len(chunks)
    )

@app.get("/train/documents", tags=["AI Training"])
async def get_documents():
    """Get list of uploaded documents with analysis data"""
    return [
        {
            "filename": doc.filename,
            "chunk_ids": doc.chunk_ids,
            "total_chunks": doc.total_chunks,
            "extractedTopics": doc.extractedTopics,
            "keyFeatures": doc.keyFeatures,
            "documentType": doc.documentType.value,
            "extractedHeadings": doc.extractedHeadings
        }
        for doc in documents_db.values()
    ]

@app.post("/train/analyze-document", response_model=DocumentAnalysis, tags=["AI Training"])
async def analyze_document(request: AnalyzeDocumentRequest):
    """Analyze document using Groq LLM"""
    print(f"Analyzing document: {request.filename}")
    
    if request.filename not in document_content:
        raise HTTPException(status_code=404, detail="Document not found")
    
    content = document_content[request.filename]
    
    system_prompt = """You are a document analysis expert. Analyze the provided document and extract key information.
    Return your response as a JSON object with the following structure:
    {
        "extractedTopics": ["topic1", "topic2", ...],
        "keyFeatures": ["feature1", "feature2", ...],
        "documentType": "proposal|technical|presentation|general",
        "extractedHeadings": [
            {
                "section": "section_name",
                "title": "heading_title", 
                "content_summary": "brief summary"
            }
        ]
    }"""
    
    prompt = f"""Analyze this document and extract key information:

Document Content:
{content[:5000]}...

Focus on:
1. Main topics and themes (extractedTopics) - maximum 5 items
2. Key features, products, or capabilities mentioned (keyFeatures) - maximum 5 items
3. Document type classification (proposal, technical, presentation, or general)
4. Section headings and their summaries - maximum 3 items

Return only the JSON response."""
    
    try:
        response = await call_groq_llm(prompt, system_prompt, max_tokens=2000)
        
        # Parse LLM response as JSON
        analysis_data = parse_json_from_llm_response(response)
        
        # Update the document in storage with analysis results
        if request.filename in documents_db:
            documents_db[request.filename].extractedTopics = analysis_data.get("extractedTopics", [])
            documents_db[request.filename].keyFeatures = analysis_data.get("keyFeatures", [])
            documents_db[request.filename].documentType = DocumentType(analysis_data.get("documentType", "general"))
            documents_db[request.filename].extractedHeadings = [
                DocumentHeading(**heading) for heading in analysis_data.get("extractedHeadings", [])
            ]
        
        print(f"Analysis complete for {request.filename}")
        return DocumentAnalysis(**analysis_data)
    
    except Exception as e:
        print(f"Analysis error: {str(e)}")
        # Fallback if analysis fails
        return DocumentAnalysis(
            extractedTopics=["Document analysis pending"],
            keyFeatures=["Features extraction pending"],
            documentType=DocumentType.general,
            extractedHeadings=[]
        )

@app.post("/train/compare-answers", response_model=CompareResponse, tags=["AI Training"])
async def compare_answers(request: CompareAnswersRequest):
    """Compare user answer with ideal answer using Groq LLM"""
    print(f"Comparing answers for question: {request.original_question[:50]}...")
    
    system_prompt = """You are an expert sales trainer and evaluator. Compare the user's answer to the ideal answer and provide detailed feedback.
    
    Return your response as a JSON object with this exact structure:
    {
        "score": 8,
        "coverage": "80%",
        "key_points_missed": ["point1", "point2"],
        "feedback": "detailed feedback text",
        "strengths": ["strength1", "strength2"],
        "improvements": ["improvement1", "improvement2"],
        "confidence_level": 85,
        "response_time": response_time_in_seconds
    }"""
    
    prompt = f"""Compare these answers for the question: "{request.original_question}"

USER ANSWER:
{request.user_answer}

IDEAL ANSWER:
{request.ideal_answer}

RESPONSE TIME: {request.response_time} seconds

Evaluate:
1. Accuracy and completeness (score 1-10)
2. Coverage percentage
3. Key points missed (maximum 3 items)
4. Detailed feedback
5. Strengths in the answer (maximum 3 items)
6. Areas for improvement (maximum 3 items)
7. Confidence level in evaluation (0-100)

Return only the JSON response."""
    
    try:
        response = await call_groq_llm(prompt, system_prompt, max_tokens=1500)
        
        # Parse LLM response as JSON
        analysis_data = parse_json_from_llm_response(response)
        analysis_data["response_time"] = request.response_time
        
        print(f"Comparison complete. Score: {analysis_data.get('score', 'N/A')}")
        return CompareResponse(**analysis_data)
    
    except Exception as e:
        print(f"Comparison error: {str(e)}")
        # Fallback analysis
        return CompareResponse(
            score=7,
            coverage="70%",
            key_points_missed=["Analysis in progress"],
            feedback="Your answer shows good understanding. The AI evaluation system is processing your response.",
            strengths=["Clear communication", "Good effort"],
            improvements=["Add more specific details", "Include concrete examples"],
            confidence_level=60,
            response_time=request.response_time
        )

@app.post("/train/generate-questions-llm", response_model=GenerateQuestionsResponse, tags=["AI Training"])
async def generate_questions_llm(request: GenerateQuestionsRequest):
    """Generate questions using Groq LLM based on documents and persona"""
    print(f"Generating questions for {request.persona} with {len(request.document_filenames)} documents")
    
    # Get content from selected documents
    selected_content = ""
    available_docs = []
    
    for filename in request.document_filenames:
        if filename in document_content:
            selected_content += f"\n\n=== {filename} ===\n{document_content[filename][:2000]}"
            available_docs.append(filename)
    
    if not selected_content:
        raise HTTPException(status_code=400, detail="No valid documents found")
    
    system_prompt = f"""You are an expert sales trainer creating questions for {request.persona} persona.
    
    Focus areas for {request.persona}:
    - Technical Lead: Architecture, integrations, technical specifications, implementation details
    - Business Manager: ROI, business value, timelines, use cases, competitive advantages  
    - C-Suite Executive: Strategic impact, scalability, risk management, long-term vision
    
    Return ONLY a JSON object with this structure:
    {{
        "questions": [
            {{
                "question": "detailed question text",
                "difficulty_level": "Easy|Medium|Hard",
                "focus_area": "specific focus area",
                "document_references": ["doc1.pdf"],
                "question_id": "unique_id"
            }}
        ]
    }}"""
    
    prompt = f"""Meeting Objective: {request.objective}

Generate {request.max_questions} questions for {request.persona} based on these documents:

{selected_content}

Create questions that:
1. Test knowledge of the document content
2. Are appropriate for {request.persona} 
3. Mix difficulty levels: Easy, Medium, Hard
4. Focus on practical scenarios they might face

Return only the JSON response with the questions array."""
    
    try:
        response = await call_groq_llm(prompt, system_prompt, max_tokens=3000)
        
        # Parse LLM response as JSON
        questions_data = parse_json_from_llm_response(response)
        
        # Ensure question_ids are unique
        for i, q in enumerate(questions_data["questions"]):
            if "question_id" not in q or not q["question_id"]:
                q["question_id"] = f"q_{uuid.uuid4().hex[:8]}"
        
        print(f"Generated {len(questions_data['questions'])} questions")
        
        return GenerateQuestionsResponse(
            questions=[GeneratedQuestion(**q) for q in questions_data["questions"]],
            status="success",
            document_count=len(available_docs),
            persona=request.persona.value
        )
    
    except Exception as e:
        print(f"Question generation error: {str(e)}")
        # Fallback questions
        fallback_questions = [
            GeneratedQuestion(
                question=f"Can you explain the key benefits of our solution for a {request.persona}?",
                difficulty_level=DifficultyLevel.medium,
                focus_area="General Benefits",
                document_references=available_docs,
                question_id=f"fallback_{uuid.uuid4().hex[:8]}"
            ),
            GeneratedQuestion(
                question=f"What concerns might a {request.persona} have about implementation?",
                difficulty_level=DifficultyLevel.medium,
                focus_area="Implementation",
                document_references=available_docs,
                question_id=f"fallback_{uuid.uuid4().hex[:8]}"
            )
        ]
        
        return GenerateQuestionsResponse(
            questions=fallback_questions,
            status="fallback",
            document_count=len(available_docs),
            persona=request.persona.value
        )

@app.post("/train/ask", response_model=IdealAnswerResponse, tags=["AI Training"])
async def ask_question(request: AskRequestTrain):
    """Get ideal answer for a question using available documents"""
    print(f"Getting ideal answer for: {request.question[:50]}...")
    
    # Find relevant content from all documents
    relevant_content = ""
    sources = []
    
    for filename, content in document_content.items():
        # Simple relevance check
        if any(word.lower() in content.lower() for word in request.question.split()):
            relevant_content += f"\n\n=== {filename} ===\n{content[:1500]}"
            sources.append(SourceTrain(filename=filename))
    
    system_prompt = """You are a knowledgeable sales assistant. Provide comprehensive answers based on the available documents."""
    
    prompt = f"""Question: {request.question}

Available context:
{relevant_content}

Provide a detailed, accurate answer based on the available information."""
    
    try:
        response = await call_groq_llm(prompt, system_prompt, max_tokens=2000)
        
        # Generate follow-up questions
        follow_ups = []
        sales_followups = []
        
        if request.include_followups:
            follow_ups = [
                "What are the implementation timelines?",
                "How does this integrate with existing systems?",
                "What are the key success metrics?"
            ]
        
        if request.include_sales_questions:
            sales_followups = [
                "What's your decision-making timeline?",
                "Who else would be involved in this decision?",
                "What's your budget range for this solution?"
            ]
        
        print("Ideal answer generated successfully")
        
        return IdealAnswerResponse(
            answer=response,
            sources=sources,
            follow_up_questions=follow_ups,
            sales_followups=sales_followups
        )
    
    except Exception as e:
        print(f"Ask question error: {str(e)}")
        return IdealAnswerResponse(
            answer="I apologize, but I'm having trouble generating an ideal answer right now. Please try again.",
            sources=sources,
            follow_up_questions=[],
            sales_followups=[]
        )
    
# --- Sentiment Analysis --- #

# Data Models
class TranscriptSegment(BaseModel):
    # ⭐ Make id optional and auto-generate if missing
    id: Optional[int] = None
    start: float
    end: float
    text: str
    language: str
    # ⭐ Make created_at optional since API returns None sometimes
    created_at: Optional[str] = None
    speaker: Optional[str] = None
    absolute_start_time: str
    absolute_end_time: str

class Transcript(BaseModel):
    id: int
    platform: str
    native_meeting_id: str
    constructed_meeting_url: str
    status: str
    start_time: str
    end_time: Optional[str]
    segments: List[TranscriptSegment]

class CriticalAlert(BaseModel):
    id: str
    timestamp: str
    phrase: str
    severity: str  # "negative_high", "negative_medium", "positive_low"
    alert_type: str  # "negative" or "positive" 
    context: str
    speaker: str
    acknowledged: bool = False
    suggestion: str = ""  # 🆕 Now shows directly in alert box, no button needed

class BuyingSignal(BaseModel):
    phrase: str
    timestamp: str
    speaker: str
    signal_type: str
    points: int
    confidence: float

class ActiveObjection(BaseModel):
    id: str
    type: str
    phrase: str
    timestamp: str
    speaker: str
    status: str  # "active" or "resolved"
    confidence: float
    context: str
class ParticipantCard(BaseModel):
    speaker: str
    role: str  # "CEO", "CTO", "Manager", "Unknown"
    status: str  # "champion", "neutral", "skeptical", "blocker"
    status_color: str  # "green", "yellow", "red"
    speaking_percentage: float
    total_words: int
    buying_signals_count: int
    concerns_count: int
    concerns_list: List[str]  # List of actual concern phrases
    action_needed: str  # Clear action for sales rep
    last_activity: str  # Timestamp of last speech

class ConversationHealth(BaseModel):
    overall_score: float  # 0-100
    health_status: str   # "healthy", "concerning", "critical"
    risk_factors: List[str]  # List of risk indicators
    positive_indicators: List[str]  # List of good signs
    recommendation: str  # What to do next
    trend: str  # "improving", "declining", "stable"

# 🆕 NEW: Engagement Score Model
class ParticipantEngagement(BaseModel):
    speaker: str
    speaking_time_seconds: float
    speaking_percentage: float
    total_words: int
    segment_count: int

# 🆕 NEW: MEDPIC Progress Model
class MedpicProgress(BaseModel):
    metrics: float  # 0-100%
    economic_buyer: float
    decision_criteria: float
    decision_process: float
    identify_pain: float
    champion: float

class SentimentAnalysisResponse(BaseModel):
    critical_alerts: List[CriticalAlert]
    buying_signals: Dict[str, Any]
    participant_cards: List[ParticipantCard]  # 🆕 NEW - replaces active_objections
    conversation_health: ConversationHealth   # 🆕 NEW - overall health score
    engagement_scores: List[ParticipantEngagement]
    medpic_progress: MedpicProgress
    last_updated: str

ROLE_KEYWORDS = {
    "ceo": ["ceo", "chief executive", "president", "founder"],
    "cto": ["cto", "chief technology", "tech lead", "technical director"],
    "cfo": ["cfo", "chief financial", "finance director", "controller"],
    "manager": ["manager", "director", "head of", "vp", "vice president"],
    "procurement": ["procurement", "purchasing", "buyer", "sourcing"],
    # Add more as needed
}

# Configuration - You can make this configurable via environment variables
CRITICAL_TRIGGERS = {
    # ⭐ RED ALERTS (Bad/Negative)
    "negative_high": [
        r"\b(end the call|not interested|waste of time|wrong solution|competitor)\b",
        r"\b(cancel|terminate|stop|abort)\b.*\b(meeting|call|demo)\b",
        r"\b(completely wrong|totally off|way off base)\b"
    ],
    "negative_medium": [
        r"\b(too expensive|can't afford|over budget|price is high)\b",
        r"\b(need to think|think about it|discuss internally|talk to team)\b",
        r"\b(not the right fit|doesn't match|not what we need)\b",
        r"\b(concerns about|worried about|hesitant)\b"
    ],
    # ⭐ GREEN ALERTS (Good/Positive)
    "positive_low": [
        r"\b(when do we start|how long does it take|what's the timeline)\b",
        r"\b(more information|additional details|want to know more)\b",
        r"\b(tell me more|sounds interesting|looks good)\b",
        r"\b(next steps|move forward|how do we proceed)\b"
    ]
}

BUYING_SIGNAL_PATTERNS = {
    "pricing": {
        "patterns": [
            r"\b(pricing for \d+|cost for \d+ users|price for \d+ licenses)\b",
            r"\b(what would it cost|how much for|pricing structure)\b",
            r"\b(budget allocation|budget planning|cost breakdown)\b"
        ],
        "points": 15
    },
    "timeline": {
        "patterns": [
            r"\b(implementation timeline|rollout schedule|go live date)\b",
            r"\b(when can we start|start date|launch timeline)\b",
            r"\b(how quickly|how fast|how soon)\b.*\b(implement|deploy|launch)\b"
        ],
        "points": 20
    },
    "decision_making": {
        "patterns": [
            r"\b(next steps|move forward|proceed with|ready to)\b",
            r"\b(sign the contract|close the deal|finalize)\b",
            r"\b(approve|approval|sign off|green light)\b"
        ],
        "points": 25
    },
    "features": {
        "patterns": [
            r"\b(love this feature|exactly what we need|perfect solution)\b",
            r"\b(this solves|addresses our|meets our needs)\b",
            r"\b(integration capabilities|API access|customization)\b"
        ],
        "points": 10
    },
    "team_expansion": {
        "patterns": [
            r"\b(add more users|expand to team|scale up)\b",
            r"\b(enterprise plan|additional licenses|more seats)\b"
        ],
        "points": 18
    }
}

OBJECTION_PATTERNS = {
    "pricing": [
        r"\b(too expensive|over budget|can't afford|price is high|costly)\b",
        r"\b(cheaper alternative|lower cost|reduce price)\b"
    ],
    "timing": [
        r"\b(not ready|too soon|need more time|busy right now)\b",
        r"\b(maybe next quarter|later this year|future consideration)\b"
    ],
    "authority": [
        r"\b(need to ask|check with|get approval|boss decides)\b",
        r"\b(not my decision|someone else decides|need buy-in)\b"
    ],
    "features": [
        r"\b(missing feature|doesn't have|lacks|not enough)\b",
        r"\b(competitor has|other solution|different approach)\b"
    ],
    "trust": [
        r"\b(not sure|skeptical|doubt|concerned|worried)\b",
        r"\b(never heard of|new company|unproven)\b"
    ]
}

# 🆕 NEW: MEDPIC Keywords Configuration
MEDPIC_KEYWORDS = {
    "metrics": [
        # Performance indicators, numbers, ROI, growth
        "revenue", "profit", "roi", "growth", "kpi", "metrics", "performance", "target", 
        "quota", "numbers", "percentage", "increase", "decrease", "improve", "efficiency",
        "productivity", "sales figures", "conversion rate", "market share", "cost savings",
        "time savings", "benchmarks", "analytics", "data", "results", "outcomes"
    ],
    "economic_buyer": [
        # Budget authority, financial decision makers
        "budget", "approve", "authorized", "decision maker", "ceo", "cfo", "vp", 
        "director", "manager", "head of", "president", "owner", "founder", "boss",
        "financial authority", "sign off", "procurement", "purchasing", "finance team",
        "budget holder", "spending authority", "final decision", "approve purchase"
    ],
    "decision_criteria": [
        # What they evaluate solutions on
        "requirements", "criteria", "features", "functionality", "must have", "nice to have",
        "evaluation", "comparison", "specifications", "capabilities", "integration",
        "security", "scalability", "reliability", "support", "training", "implementation",
        "user friendly", "ease of use", "compatibility", "performance", "speed"
    ],
    "decision_process": [
        # How they make decisions, timeline, steps
        "process", "timeline", "steps", "phases", "approval process", "evaluation process",
        "next steps", "when", "how long", "duration", "schedule", "plan", "roadmap",
        "committee", "team decision", "consensus", "vote", "review process", "trial",
        "pilot", "demo", "proof of concept", "implementation plan"
    ],
    "identify_pain": [
        # Problems, challenges, pain points
        "problem", "challenge", "issue", "pain", "frustration", "difficulty", "struggle",
        "bottleneck", "inefficiency", "waste", "manual", "time consuming", "error prone",
        "costly", "expensive", "slow", "complicated", "difficult", "annoying",
        "broken", "not working", "failing", "limitation", "constraint", "obstacle"
    ],
    "champion": [
        # Internal advocates, supporters
        "champion", "advocate", "supporter", "internal sponsor", "ally", "recommend",
        "endorse", "promote", "support", "backing", "help us", "work with us",
        "collaborate", "partnership", "relationship", "trust", "believe in",
        "convinced", "sold on", "excited about", "enthusiastic", "positive"
    ]
}

POSITIVE_WORDS = [
    "excellent", "perfect", "amazing", "love", "great", "fantastic", "wonderful",
    "exactly", "definitely", "absolutely", "yes", "agree", "correct", "right",
    "good", "nice", "awesome", "outstanding", "impressive", "solid", "strong"
]

NEGATIVE_WORDS = [
    "terrible", "awful", "hate", "horrible", "bad", "wrong", "no", "never",
    "impossible", "difficult", "problem", "issue", "concern", "worry", "doubt",
    "confused", "unclear", "complicated", "expensive", "costly", "slow"
]

ALERT_SUGGESTIONS = {
    "negative_high": [
        "💡 **Pause and Listen**: Let them express their concerns fully before responding",
        "🔄 **Redirect Focus**: Ask what specific aspects aren't working for them",
        "🤝 **Find Common Ground**: Identify shared goals or pain points",
        "⏰ **Schedule Follow-up**: Suggest continuing the conversation when they're ready"
    ],
    "negative_medium": [
        "💰 **Address Budget**: Present ROI calculations and payment options",
        "👥 **Involve Decision Makers**: Ask who else should be part of the conversation",
        "📊 **Show Value**: Provide case studies and success stories",
        "🎯 **Customize Solution**: Adjust the proposal to fit their specific needs"
    ],
    "positive_low": [
        "🚀 **Seize Momentum**: Provide detailed implementation timeline",
        "📋 **Share Resources**: Send relevant documentation and materials",
        "📅 **Schedule Next Steps**: Book follow-up meetings or demos",
        "🎯 **Close on Interest**: Ask for specific commitments or next actions"
    ]
}

class SentimentAnalyzer:
    def __init__(self):
        self.alert_counter = 0
        self.objection_counter = 0
        self.segment_id_counter = 0
        self.recent_alerts = []
        self.participant_history = {}  # Track participant data over time
        
    def _is_duplicate_alert(self, phrase: str, speaker: str, timestamp: str) -> bool:
        """Check if similar alert was sent in the past 3 minutes"""
        try:
            current_time = datetime.fromisoformat(timestamp.replace('Z', '+00:00'))
            three_minutes_ago = current_time - timedelta(minutes=3)
            
            # Clean up old alerts (older than 3 minutes)
            self.recent_alerts = [
                alert for alert in self.recent_alerts 
                if datetime.fromisoformat(alert['timestamp'].replace('Z', '+00:00')) > three_minutes_ago
            ]
            
            # Check for duplicates
            for recent_alert in self.recent_alerts:
                if (recent_alert['phrase'].lower() == phrase.lower() and 
                    recent_alert['speaker'] == speaker):
                    return True
            
            return False
            
        except ValueError:
            # If timestamp parsing fails, use simple approach
            return any(
                alert['phrase'].lower() == phrase.lower() and alert['speaker'] == speaker 
                for alert in self.recent_alerts[-10:]  # Check last 10 alerts as fallback
            )
        
    async def get_transcript_data(self, platform: str, meeting_id: str) -> Transcript:
        """Fetch transcript data from the existing endpoint"""
        try:
            async with httpx.AsyncClient() as client:
                response = await client.get(f"{BASE_URL}/vexa/transcripts/{platform}/{meeting_id}")
                if response.status_code == 200:
                    data = response.json()
                    
                    # ⭐ Fix segments that don't have IDs or created_at
                    if 'segments' in data:
                        for i, segment in enumerate(data['segments']):
                            if 'id' not in segment or segment['id'] is None:
                                segment['id'] = i + 1  # Auto-generate ID
                            if 'created_at' not in segment or segment['created_at'] is None:
                                segment['created_at'] = segment.get('absolute_start_time', '')
                    
                    return Transcript(**data)
                else:
                    logger.error(f"Failed to fetch transcript: {response.status_code}")
                    raise HTTPException(status_code=response.status_code, detail="Failed to fetch transcript")
        except Exception as e:
            logger.error(f"Error fetching transcript: {str(e)}")
            raise HTTPException(status_code=500, detail=f"Error fetching transcript: {str(e)}")

    def analyze_critical_alerts(self, segments: List[TranscriptSegment]) -> List[CriticalAlert]:
        """Analyze transcript for critical alert triggers"""
        alerts = []
        
        for segment in segments:
            text = segment.text.lower()
            
            # Skip AI bot messages
            if segment.speaker and 'spiked' in segment.speaker.lower():
                continue
                
            for severity, patterns in CRITICAL_TRIGGERS.items():
                for pattern in patterns:
                    matches = re.finditer(pattern, text, re.IGNORECASE)
                    for match in matches:
                        phrase = match.group()
                        speaker = segment.speaker or "Unknown"
                        timestamp = segment.absolute_start_time
                        
                        # 🆕 NEW: Check for duplicate alerts in past 3 minutes
                        if self._is_duplicate_alert(phrase, speaker, timestamp):
                            logger.info(f"Skipping duplicate alert: '{phrase}' from {speaker}")
                            continue
                        
                        self.alert_counter += 1
                        
                        # Determine alert type and get suggestion
                        if "negative" in severity:
                            alert_type = "negative"
                        else:
                            alert_type = "positive"
                        
                        # Get random suggestion for this severity level
                        import random
                        suggestions = ALERT_SUGGESTIONS.get(severity, ["Consider your response carefully"])
                        suggestion = random.choice(suggestions)
                        
                        # Get context (surrounding text)
                        start_idx = max(0, match.start() - 50)
                        end_idx = min(len(text), match.end() + 50)
                        context = text[start_idx:end_idx]
                        
                        alert = CriticalAlert(
                            id=f"alert_{self.alert_counter}",
                            timestamp=timestamp,
                            phrase=phrase,
                            severity=severity,
                            alert_type=alert_type,
                            context=context,
                            speaker=speaker,
                            suggestion=suggestion  # 🆕 NEW: Suggestion now included directly
                        )
                        alerts.append(alert)
                        
                        # 🆕 NEW: Add to recent alerts tracking
                        self.recent_alerts.append({
                            'phrase': phrase,
                            'speaker': speaker,
                            'timestamp': timestamp
                        })
        
        # Remove duplicates and sort by timestamp
        unique_alerts = []
        seen_phrases = set()
        for alert in alerts:
            key = f"{alert.phrase}_{alert.timestamp}_{alert.speaker}"
            if key not in seen_phrases:
                seen_phrases.add(key)
                unique_alerts.append(alert)
                
        return sorted(unique_alerts, key=lambda x: x.timestamp, reverse=True)

    def analyze_buying_signals(self, segments: List[TranscriptSegment]) -> Dict[str, Any]:
        """Analyze transcript for buying signals"""
        signals = []
        signals_by_type = defaultdict(int)
        total_score = 0
        
        for segment in segments:
            text = segment.text.lower()
            
            # Skip AI bot messages
            if segment.speaker and 'spiked' in segment.speaker.lower():
                continue
                
            for signal_type, config in BUYING_SIGNAL_PATTERNS.items():
                for pattern in config["patterns"]:
                    matches = re.finditer(pattern, text, re.IGNORECASE)
                    for match in matches:
                        confidence = 0.8  # Base confidence
                        
                        # Increase confidence if multiple buying signals in same segment
                        if len(list(re.finditer(r'\b(pricing|timeline|next steps|implement|start)\b', text))) > 1:
                            confidence = min(0.95, confidence + 0.1)
                        
                        signal = BuyingSignal(
                            phrase=match.group(),
                            timestamp=segment.absolute_start_time,
                            speaker=segment.speaker or "Unknown",
                            signal_type=signal_type,
                            points=config["points"],
                            confidence=confidence
                        )
                        signals.append(signal)
                        signals_by_type[signal_type] += config["points"]
                        total_score += config["points"]
        
        # ⭐ FIX: Use timezone-aware datetime
        now = datetime.now(timezone.utc)
        last_10_min = now - timedelta(minutes=10)
        prev_10_min = now - timedelta(minutes=20)
        
        try:
            recent_signals = [s for s in signals if datetime.fromisoformat(s.timestamp.replace('Z', '+00:00')) > last_10_min]
            previous_signals = [s for s in signals if prev_10_min < datetime.fromisoformat(s.timestamp.replace('Z', '+00:00')) <= last_10_min]
        except ValueError:
            # Fallback if timestamp parsing fails
            recent_signals = signals[-5:] if len(signals) > 5 else signals
            previous_signals = signals[-10:-5] if len(signals) > 10 else []
        
        recent_score = sum(s.points for s in recent_signals)
        previous_score = sum(s.points for s in previous_signals)
        
        if previous_score == 0:
            trend = "neutral" if recent_score == 0 else "increasing"
        else:
            trend = "increasing" if recent_score > previous_score else "decreasing" if recent_score < previous_score else "stable"
        
        return {
            "total_score": total_score,
            "signals_by_type": dict(signals_by_type),
            "trend": trend,
            "recent_signals": [s.dict() for s in signals[-5:]],  # Last 5 signals
            "signal_count": len(signals)
        }

    def detect_participant_role(self, speaker: str, segments: List[TranscriptSegment]) -> str:
        """Detect participant role based on name and speech patterns"""
        if not speaker or speaker == "Unknown":
            return "Unknown"
        
        speaker_lower = speaker.lower()
        
        # Check for role keywords in speaker name or their speech
        speaker_segments = [s for s in segments if s.speaker == speaker]
        combined_text = " ".join([s.text.lower() for s in speaker_segments])
        
        for role, keywords in ROLE_KEYWORDS.items():
            if any(keyword in speaker_lower for keyword in keywords):
                return role.upper()
            if any(keyword in combined_text for keyword in keywords):
                return role.upper()
        
        # Default role detection based on speech patterns
        if len(speaker_segments) > 20:  # Very active speaker
            return "DECISION_MAKER"
        elif any(word in combined_text for word in ["technical", "integration", "system", "api"]):
            return "TECHNICAL"
        elif any(word in combined_text for word in ["budget", "cost", "price", "financial"]):
            return "FINANCIAL"
        else:
            return "STAKEHOLDER"

    def analyze_participant_sentiment(self, speaker: str, segments: List[TranscriptSegment]) -> Dict[str, Any]:
        """Analyze individual participant sentiment and concerns"""
        speaker_segments = [s for s in segments if s.speaker == speaker]
        
        if not speaker_segments:
            return {
                "sentiment_score": 0,
                "buying_signals": 0,
                "concerns": [],
                "last_activity": "",
                "dominant_sentiment": "neutral"
            }
        
        # Calculate sentiment metrics
        buying_signals = 0
        concerns = []
        sentiment_score = 0
        
        for segment in speaker_segments:
            text = segment.text.lower()
            
            # Check for buying signals
            for signal_type, config in BUYING_SIGNAL_PATTERNS.items():
                for pattern in config["patterns"]:
                    if re.search(pattern, text, re.IGNORECASE):
                        buying_signals += 1
            
            # Check for concerns/objections
            for objection_type, patterns in OBJECTION_PATTERNS.items():
                for pattern in patterns:
                    matches = re.findall(pattern, text, re.IGNORECASE)
                    for match in matches:
                        concerns.append(match)
            
            # Simple sentiment scoring
            positive_words = sum(1 for word in POSITIVE_WORDS if word in text)
            negative_words = sum(1 for word in NEGATIVE_WORDS if word in text)
            sentiment_score += (positive_words - negative_words)
        
        # Determine dominant sentiment
        if sentiment_score > 2:
            dominant_sentiment = "positive"
        elif sentiment_score < -2:
            dominant_sentiment = "negative"
        else:
            dominant_sentiment = "neutral"
        
        return {
            "sentiment_score": sentiment_score,
            "buying_signals": buying_signals,
            "concerns": list(set(concerns))[:3],  # Top 3 unique concerns
            "last_activity": speaker_segments[-1].absolute_start_time if speaker_segments else "",
            "dominant_sentiment": dominant_sentiment
        }

    def analyze_participant_cards(self, segments: List[TranscriptSegment]) -> List[ParticipantCard]:
        """Generate participant status cards"""
        # Get engagement scores first
        engagement_scores = self.analyze_engagement_scores(segments)
        
        participant_cards = []
        
        for engagement in engagement_scores:
            speaker = engagement.speaker
            
            # Skip AI bot
            if 'spiked' in speaker.lower():
                continue
            
            # Detect role
            role = self.detect_participant_role(speaker, segments)
            
            # Analyze sentiment
            sentiment_data = self.analyze_participant_sentiment(speaker, segments)
            
            # Determine status and color
            concerns_count = len(sentiment_data["concerns"])
            buying_signals_count = sentiment_data["buying_signals"]
            speaking_percentage = engagement.speaking_percentage
            
            # Status logic
            if buying_signals_count >= 2 and concerns_count == 0 and speaking_percentage > 15:
                status = "champion"
                status_color = "green"
                action_needed = "Keep engaging - they're your advocate!"
            elif concerns_count >= 2 or speaking_percentage < 10:
                status = "blocker" if concerns_count >= 3 else "skeptical"
                status_color = "red"
                if speaking_percentage < 10:
                    action_needed = f"Get them talking - only {speaking_percentage:.1f}% engagement"
                else:
                    action_needed = f"Address their concerns: {', '.join(sentiment_data['concerns'][:2])}"
            elif concerns_count == 1 or (speaking_percentage < 20 and buying_signals_count == 0):
                status = "neutral"
                status_color = "yellow"
                action_needed = "Need to engage more - unclear position"
            else:
                status = "positive"
                status_color = "green"
                action_needed = "Maintain momentum"
            
            participant_card = ParticipantCard(
                speaker=speaker,
                role=role,
                status=status,
                status_color=status_color,
                speaking_percentage=speaking_percentage,
                total_words=engagement.total_words,
                buying_signals_count=buying_signals_count,
                concerns_count=concerns_count,
                concerns_list=sentiment_data["concerns"],
                action_needed=action_needed,
                last_activity=sentiment_data["last_activity"]
            )
            
            participant_cards.append(participant_card)
        
        # Sort by risk level (red first, then yellow, then green)
        color_priority = {"red": 0, "yellow": 1, "green": 2}
        participant_cards.sort(key=lambda x: color_priority.get(x.status_color, 3))
        
        return participant_cards

    def analyze_conversation_health(self, segments: List[TranscriptSegment], participant_cards: List[ParticipantCard]) -> ConversationHealth:
        """Analyze overall conversation health"""
        if not segments:
            return ConversationHealth(
                overall_score=50.0,
                health_status="neutral",
                risk_factors=[],
                positive_indicators=[],
                recommendation="Start the conversation",
                trend="stable"
            )
        
        # Calculate health metrics
        risk_factors = []
        positive_indicators = []
        score = 50.0  # Start neutral
        
        # Analyze participant distribution
        red_participants = [p for p in participant_cards if p.status_color == "red"]
        green_participants = [p for p in participant_cards if p.status_color == "green"]
        
        # Risk factors
        if len(red_participants) > len(green_participants):
            risk_factors.append("More skeptical than positive participants")
            score -= 15
        
        if any(p.speaking_percentage < 5 for p in participant_cards):
            risk_factors.append("Key participants not engaging")
            score -= 10
        
        total_concerns = sum(p.concerns_count for p in participant_cards)
        if total_concerns > 3:
            risk_factors.append(f"{total_concerns} concerns raised")
            score -= 10
        
        # Check for decision maker engagement
        decision_makers = [p for p in participant_cards if p.role in ["CEO", "DECISION_MAKER", "MANAGER"]]
        if decision_makers and all(p.speaking_percentage < 15 for p in decision_makers):
            risk_factors.append("Decision makers not actively participating")
            score -= 20
        
        # Positive indicators
        if len(green_participants) > 1:
            positive_indicators.append("Multiple champions identified")
            score += 15
        
        total_buying_signals = sum(p.buying_signals_count for p in participant_cards)
        if total_buying_signals > 2:
            positive_indicators.append(f"{total_buying_signals} buying signals detected")
            score += 10
        
        if any(p.speaking_percentage > 30 for p in participant_cards if p.status_color == "green"):
            positive_indicators.append("Champions are actively engaged")
            score += 10
        
        # Determine health status
        if score >= 70:
            health_status = "healthy"
            recommendation = "Deal progressing well - maintain momentum"
        elif score >= 40:
            health_status = "concerning"
            recommendation = "Address key concerns and engage quiet participants"
        else:
            health_status = "critical"
            recommendation = "Deal at risk - immediate intervention needed"
        
        # Simple trend calculation (would be more sophisticated with historical data)
        recent_segments = segments[-10:] if len(segments) > 10 else segments
        recent_positive = sum(1 for s in recent_segments if any(word in s.text.lower() for word in POSITIVE_WORDS))
        recent_negative = sum(1 for s in recent_segments if any(word in s.text.lower() for word in NEGATIVE_WORDS))
        
        if recent_positive > recent_negative:
            trend = "improving"
        elif recent_negative > recent_positive:
            trend = "declining"
        else:
            trend = "stable"
        
        return ConversationHealth(
            overall_score=round(max(0, min(100, score)), 1),
            health_status=health_status,
            risk_factors=risk_factors,
            positive_indicators=positive_indicators,
            recommendation=recommendation,
            trend=trend
        )

    # 🆕 NEW: Engagement Score Analysis
    def analyze_engagement_scores(self, segments: List[TranscriptSegment]) -> List[ParticipantEngagement]:
        """Analyze speaking time percentage for each participant"""
        if not segments:
            return []
        
        speaker_stats = defaultdict(lambda: {
            'total_time': 0.0,
            'word_count': 0,
            'segment_count': 0
        })
        
        total_meeting_time = 0.0
        
        for segment in segments:
            # Skip AI bot messages
            if segment.speaker and 'spiked' in segment.speaker.lower():
                continue
                
            speaker = segment.speaker or "Unknown"
            segment_duration = segment.end - segment.start
            word_count = len(segment.text.split())
            
            speaker_stats[speaker]['total_time'] += segment_duration
            speaker_stats[speaker]['word_count'] += word_count
            speaker_stats[speaker]['segment_count'] += 1
            total_meeting_time += segment_duration
        
        # Calculate percentages
        engagement_scores = []
        for speaker, stats in speaker_stats.items():
            if total_meeting_time > 0:
                speaking_percentage = (stats['total_time'] / total_meeting_time) * 100
            else:
                speaking_percentage = 0.0
                
            engagement = ParticipantEngagement(
                speaker=speaker,
                speaking_time_seconds=round(stats['total_time'], 2),
                speaking_percentage=round(speaking_percentage, 1),
                total_words=stats['word_count'],
                segment_count=stats['segment_count']
            )
            engagement_scores.append(engagement)
        
        # Sort by speaking percentage (highest first)
        return sorted(engagement_scores, key=lambda x: x.speaking_percentage, reverse=True)

    # 🆕 NEW: MEDPIC Progress Analysis
    def analyze_medpic_progress(self, segments: List[TranscriptSegment]) -> MedpicProgress:
        """Analyze MEDPIC coverage progress in the meeting"""
        if not segments:
            return MedpicProgress(
                metrics=0.0,
                economic_buyer=0.0,
                decision_criteria=0.0,
                decision_process=0.0,
                identify_pain=0.0,
                champion=0.0
            )
        
        # Count mentions of each MEDPIC category
        medpic_mentions = defaultdict(int)
        total_medpic_mentions = 0
        
        for segment in segments:
            # Skip AI bot messages
            if segment.speaker and 'spiked' in segment.speaker.lower():
                continue
                
            text = segment.text.lower()
            
            # Count mentions for each MEDPIC category
            for category, keywords in MEDPIC_KEYWORDS.items():
                for keyword in keywords:
                    if keyword in text:
                        medpic_mentions[category] += text.count(keyword)
                        total_medpic_mentions += text.count(keyword)
        
        # Calculate progress percentages using the formula:
        # category_percentage = (mentions_of_category / total_medpic_mentions) * 100
        # Then normalize to 0-100% based on relative coverage
        
        if total_medpic_mentions == 0:
            # No MEDPIC topics discussed yet - all red (0%)
            return MedpicProgress(
                metrics=0.0,
                economic_buyer=0.0,
                decision_criteria=0.0,
                decision_process=0.0,
                identify_pain=0.0,
                champion=0.0
            )
        
        # Calculate raw percentages
        raw_percentages = {}
        for category in MEDPIC_KEYWORDS.keys():
            if total_medpic_mentions > 0:
                raw_percentages[category] = (medpic_mentions[category] / total_medpic_mentions) * 100
            else:
                raw_percentages[category] = 0.0
        
        # Normalize to 0-100% scale (if any category is mentioned, it gets at least some progress)
        # We'll use a logarithmic scale to show progress more meaningfully
        normalized_percentages = {}
        for category, raw_pct in raw_percentages.items():
            if raw_pct == 0:
                normalized_percentages[category] = 0.0
            else:
                # Convert to a more meaningful scale where having mentions = some progress
                # Minimum 10% if mentioned, scaling up to 100% based on relative frequency
                normalized_percentages[category] = min(100.0, max(10.0, raw_pct * 5))
        
        return MedpicProgress(
            metrics=round(normalized_percentages.get('metrics', 0.0), 1),
            economic_buyer=round(normalized_percentages.get('economic_buyer', 0.0), 1),
            decision_criteria=round(normalized_percentages.get('decision_criteria', 0.0), 1),
            decision_process=round(normalized_percentages.get('decision_process', 0.0), 1),
            identify_pain=round(normalized_percentages.get('identify_pain', 0.0), 1),
            champion=round(normalized_percentages.get('champion', 0.0), 1)
        )


# Initialize analyzer
analyzer = SentimentAnalyzer()
acknowledged_alerts = set()

@app.get("/sentiment/analysis/{platform}/{meeting_id}", response_model=SentimentAnalysisResponse, tags=["Sentiment Analysis"])
async def get_sentiment_analysis(platform: str, meeting_id: str):
    """
    Get comprehensive sentiment analysis for a meeting
    
    This endpoint analyzes the real-time transcript and provides:
    - Critical alerts for concerning phrases
    - Buying signals with scoring
    - Participant status cards with individual sentiment
    - Overall conversation health score
    - Engagement scores showing % speaking time per participant
    - MEDPIC progress tracking with progress bars
    """
    try:
        # Get transcript data
        transcript_data = await analyzer.get_transcript_data(platform, meeting_id)
        
        if not transcript_data.segments:
            logger.warning(f"No transcript segments found for {platform}/{meeting_id}")
            return SentimentAnalysisResponse(
                critical_alerts=[],
                buying_signals={"total_score": 0, "signals_by_type": {}, "trend": "neutral", "recent_signals": [], "signal_count": 0},
                participant_cards=[],
                conversation_health=ConversationHealth(
                    overall_score=50.0, health_status="neutral", risk_factors=[], 
                    positive_indicators=[], recommendation="Start the conversation", trend="stable"
                ),
                engagement_scores=[],
                medpic_progress=MedpicProgress(
                    metrics=0.0, economic_buyer=0.0, decision_criteria=0.0,
                    decision_process=0.0, identify_pain=0.0, champion=0.0
                ),
                last_updated=datetime.now().isoformat()
            )
        
        # Run all analyses
        logger.info(f"Analyzing {len(transcript_data.segments)} segments for {platform}/{meeting_id}")
        
        critical_alerts = analyzer.analyze_critical_alerts(transcript_data.segments)
        buying_signals = analyzer.analyze_buying_signals(transcript_data.segments)
        participant_cards = analyzer.analyze_participant_cards(transcript_data.segments)  # 🆕 NEW
        conversation_health = analyzer.analyze_conversation_health(transcript_data.segments, participant_cards)  # 🆕 NEW
        engagement_scores = analyzer.analyze_engagement_scores(transcript_data.segments)
        medpic_progress = analyzer.analyze_medpic_progress(transcript_data.segments)
        
        logger.info(f"Analysis complete: {len(critical_alerts)} alerts, {buying_signals['signal_count']} signals, {len(participant_cards)} participants")
        
        return SentimentAnalysisResponse(
            critical_alerts=critical_alerts,
            buying_signals=buying_signals,
            participant_cards=participant_cards,      # 🆕 NEW
            conversation_health=conversation_health,  # 🆕 NEW
            engagement_scores=engagement_scores,
            medpic_progress=medpic_progress,
            last_updated=datetime.now().isoformat()
        )
        
    except HTTPException:
        raise
    except Exception as e:
        logger.error(f"Error in sentiment analysis: {str(e)}")
        raise HTTPException(status_code=500, detail=f"Internal error: {str(e)}")

# 🆕 NEW INDIVIDUAL ENDPOINTS
@app.get("/sentiment/participants/{platform}/{meeting_id}", tags=["Sentiment Analysis"])
async def get_participant_cards_only(platform: str, meeting_id: str):
    """Get only participant status cards"""
    transcript_data = await analyzer.get_transcript_data(platform, meeting_id)
    participant_cards = analyzer.analyze_participant_cards(transcript_data.segments)
    return {"participant_cards": participant_cards}

@app.get("/sentiment/health/{platform}/{meeting_id}", tags=["Sentiment Analysis"])
async def get_conversation_health_only(platform: str, meeting_id: str):
    """Get only conversation health score"""
    transcript_data = await analyzer.get_transcript_data(platform, meeting_id)
    participant_cards = analyzer.analyze_participant_cards(transcript_data.segments)
    health = analyzer.analyze_conversation_health(transcript_data.segments, participant_cards)
    return {"conversation_health": health}

# Optional: Additional endpoints for individual features
@app.get("/sentiment/alerts/{platform}/{meeting_id}", tags=["Sentiment Analysis"])
async def get_critical_alerts_only(platform: str, meeting_id: str):
    """Get only critical alerts"""
    transcript_data = await analyzer.get_transcript_data(platform, meeting_id)
    alerts = analyzer.analyze_critical_alerts(transcript_data.segments)
    return {"critical_alerts": alerts}

@app.get("/sentiment/engagement/{platform}/{meeting_id}", tags=["Sentiment Analysis"])
async def get_engagement_scores_only(platform: str, meeting_id: str):
    """Get only engagement scores"""
    transcript_data = await analyzer.get_transcript_data(platform, meeting_id)
    engagement = analyzer.analyze_engagement_scores(transcript_data.segments)
    return {"engagement_scores": engagement}

@app.get("/sentiment/medpic/{platform}/{meeting_id}", tags=["Sentiment Analysis"])
async def get_medpic_progress_only(platform: str, meeting_id: str):
    """Get only MEDPIC progress"""
    transcript_data = await analyzer.get_transcript_data(platform, meeting_id)
    medpic = analyzer.analyze_medpic_progress(transcript_data.segments)
    return {"medpic_progress": medpic}

@app.post("/sentiment/alerts/{alert_id}/acknowledge", tags=["Sentiment Analysis"])
async def acknowledge_alert(alert_id: str):
    """Marks a critical alert as acknowledged by a user."""
    if not isinstance(alert_id, str) or not alert_id:
        raise HTTPException(status_code=400, detail="A valid alert_id string is required.")
    acknowledged_alerts.add(alert_id)
    logger.info(f"Alert '{alert_id}' has been acknowledged.")
    return {"status": "acknowledged", "alert_id": alert_id}

@app.get("/sentiment/alerts/{alert_id}/suggestion", tags=["Sentiment Analysis"])
async def get_alert_suggestion(alert_id: str):
    """Get AI suggestion for handling an alert - DEPRECATED: Suggestions now included in alert directly"""
    return {
        "alert_id": alert_id,
        "message": "Suggestions are now included directly in the alert object. Check the 'suggestion' field.",
        "suggestions": [
            "This endpoint is deprecated - suggestions are now in the alert.suggestion field"
        ]
    }

# --- Graph RAG --- #

# API Keys
COHERE_API_KEY="Irw89AunBE04841KAEnLp3TJH1Ubyrk72Uhe5qZY"

# Neo4j Database
NEO4J_URI="neo4j+s://e2885109.databases.neo4j.io"
NEO4J_USER="neo4j"
NEO4J_PASSWORD="aeGjnEVaBsRR3v78ZaWsvNFNjjoCqaEr8RO40o7XiWs"

from langchain_groq import ChatGroq
from langchain_openai import OpenAIEmbeddings
from langchain.schema import HumanMessage
from neo4j import GraphDatabase
import cohere

model_llama = ChatGroq(
    model=GROQ_MODEL, 
    groq_api_key=GROQ_API_KEY,
    streaming=True
)
embeddings = OpenAIEmbeddings(
    model="text-embedding-3-small",
    openai_api_key=OPENAI_API_KEY
)
co = cohere.Client(api_key=COHERE_API_KEY)

# Neo4j connection
driver = GraphDatabase.driver(
    NEO4J_URI, 
    auth=(NEO4J_USER, NEO4J_PASSWORD)
)

class QueryRequest(BaseModel):
    query: str

def retrieve_docs(query):
    """Retrieve documents from Neo4j using vector similarity"""
    query_embed = embeddings.embed_query(query)
    cypher_query = """
    MATCH (n:Chunk)
    WHERE n.embedding is NOT NULL
    WITH n, gds.similarity.cosine(n.embedding, $query_vector) AS score
    RETURN n, score
    ORDER BY score DESC
    limit 20
    """
    with driver.session() as session:
        result = session.run(cypher_query, query_vector=query_embed)
        results = result.data()
    return results

def rerank_docs(query, results_cypher):
    """Rerank documents using Cohere"""
    docs = []
    for doc in results_cypher:
        docs.append(doc['n']['text'])
    
    ranked_res = co.rerank(model='rerank-v3.5', query=query, documents=docs, top_n=7)
    final_docs = []
    for res in ranked_res.results:
        index = res.index
        final_docs.append(docs[index])
    return final_docs

@app.post("/graph/ask", summary="Get answers via Graph RAG (Testing)", tags=["External AI"])
async def stream_rag_response(request: QueryRequest):
    async def generate():
        # Retrieve documents
        results = retrieve_docs(request.query)
        
        
        reranked_docs = rerank_docs(request.query, results)
        
       
        context = '\n\n'.join(reranked_docs)
        prompt = f"""You are the rag agent for company KoreAI. Your application is customer facing. You are answering questions from the knowledge
      base of all of company's data. 
      Given the user query and the context answer the question\nquery: {request.query}\n\n
      ## context: \n {context} \n\n\n\n
      ## Important
      ### Give your answer in markdown format. give tables wherever necessary. \n
      ### highlight keywords(as markdown bold) that are important and might help in answering the question
      ### If the question is not relevant to sales or koreai or related subjects, don't answer.
      
      """
        
       
        messages = [HumanMessage(content=prompt)]
        async for chunk in model_llama.astream(messages):
            if chunk.content:
                yield chunk.content
    
    return StreamingResponse(generate(), media_type="text/plain")

# --- Ask Beyond Documents --- #

OPENAI_CHAT_MODEL = "gpt-4o"

class AskResponse(BaseModel):
    answer: str
    question: str
    sources: List[dict]
    client_followups: List[str] = Field(default_factory=list)
    sales_followups: List[str] = Field(default_factory=list)

class AskBeyondResponse(BaseModel):
    answer: str
    question: str

def _initialize_client(client_name: str, init_func, *args, **kwargs):
    """Generic client initializer to reduce code duplication."""
    logger.info(f"Initializing {client_name} client...")
    try:
        client = init_func(*args, **kwargs)
        logger.info(f"{client_name} client initialized successfully.")
        return client
    except Exception as e:
        logger.critical(f"FATAL: Failed to initialize {client_name} client: {e}", exc_info=True)
        # This will stop the application startup if a critical client fails.
        raise RuntimeError(f"Could not initialize {client_name}: {e}") from e

def get_openai_client() -> openai.OpenAI:
    global openai_client
    if openai_client is None:
        if not OPENAI_API_KEY:
            raise RuntimeError("OPENAI_API_KEY is not configured.")
        openai_client = _initialize_client("OpenAI", openai.OpenAI, api_key=OPENAI_API_KEY)
    return openai_client

@app.post("/askBeyond", response_model=AskBeyondResponse, summary="Ask a question with web search capabilities", tags=["External AI"])
async def ask_beyond(request: AskRequest):
    """Answers a question using a powerful external AI model with web search."""
    question = request.question.strip()
    if not question:
        raise HTTPException(status_code=400, detail="Question cannot be empty.")
    
    client = get_openai_client()
    try:
        completion = await asyncio.to_thread(
            client.chat.completions.create,
            model=OPENAI_CHAT_MODEL,
            messages=[
                {"role": "system", "content": "You are a helpful assistant with access to real-time information. Provide clear, concise, and accurate answers."},
                {"role": "user", "content": question}
            ]
        )
        answer = completion.choices[0].message.content
        if not answer:
             raise HTTPException(status_code=500, detail="AI model returned an empty answer.")
        return AskBeyondResponse(answer=answer.strip(), question=question)
    except openai.APIError as e:
        logger.error(f"OpenAI API call failed: {e}", exc_info=True)
        raise HTTPException(status_code=502, detail=f"Failed to get a response from the AI model: {e}")
    except Exception as e:
        logger.error(f"An unexpected error occurred in ask_beyond: {e}", exc_info=True)
        raise HTTPException(status_code=500, detail="An internal server error occurred.")

