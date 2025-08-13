import io
import hashlib
import logging
from typing import Dict, Any, List
from pathlib import Path
from datetime import datetime

import fitz  # PyMuPDF
from docx import Document as DocxDocument
import openpyxl
from pptx import Presentation
import PyPDF2

from sentence_transformers import SentenceTransformer
import chromadb
from chromadb.config import Settings

logger = logging.getLogger(__name__)

class EnhancedDocumentProcessor:
    def __init__(self, embedding_model: SentenceTransformer, collection):
        self.embedding_model = embedding_model
        self.collection = collection
        self.supported_formats = ['.pdf', '.docx', '.txt', '.md', '.xlsx', '.pptx']
    
    async def process_document(self, file_content: bytes, filename: str, content_type: str) -> Dict[str, Any]:
        """Process uploaded document and extract content"""
        try:
            file_extension = Path(filename).suffix.lower()
            
            if file_extension not in self.supported_formats:
                raise ValueError(f"Unsupported file format: {file_extension}")
            
            # Extract text based on file type
            extracted_text = await self._extract_text(file_content, file_extension)
            
            if not extracted_text or len(extracted_text.strip()) < 50:
                raise ValueError("Document contains insufficient text content")
            
            # Create document ID
            document_id = hashlib.md5(f"{filename}_{len(file_content)}".encode()).hexdigest()
            
            # Create embeddings and store in vector DB
            embeddings_count = await self._store_document_embeddings(extracted_text, document_id, filename)
            
            return {
                'filename': filename,
                'content_type': content_type,
                'size': len(file_content),
                'extracted_text': extracted_text[:500] + "..." if len(extracted_text) > 500 else extracted_text,
                'embeddings_count': embeddings_count,
                'document_id': document_id,
                'full_text_length': len(extracted_text)
            }
            
        except Exception as e:
            logger.error(f"Document processing failed: {e}")
            raise
    
    async def _extract_text(self, content: bytes, file_extension: str) -> str:
        """Extract text from different file formats"""
        try:
            if file_extension == '.pdf':
                return self._extract_pdf_text(content)
            elif file_extension == '.docx':
                return self._extract_docx_text(content)
            elif file_extension == '.xlsx':
                return self._extract_xlsx_text(content)
            elif file_extension == '.pptx':
                return self._extract_pptx_text(content)
            elif file_extension in ['.txt', '.md']:
                return content.decode('utf-8', errors='ignore')
            else:
                raise ValueError(f"Unsupported file format: {file_extension}")
        except Exception as e:
            logger.error(f"Text extraction failed: {e}")
            raise
    
    def _extract_pdf_text(self, content: bytes) -> str:
        """Extract text from PDF using PyMuPDF"""
        try:
            doc = fitz.open(stream=content, filetype="pdf")
            text = ""
            for page in doc:
                text += page.get_text()
            doc.close()
            return text
        except Exception as e:
            logger.error(f"PDF extraction failed: {e}")
            # Fallback to PyPDF2
            try:
                pdf_reader = PyPDF2.PdfReader(io.BytesIO(content))
                text = ""
                for page in pdf_reader.pages:
                    text += page.extract_text() + "\n"
                return text
            except Exception as e2:
                logger.error(f"PyPDF2 fallback also failed: {e2}")
                raise
    
    def _extract_docx_text(self, content: bytes) -> str:
        """Extract text from DOCX"""
        try:
            doc = DocxDocument(io.BytesIO(content))
            text = ""
            for paragraph in doc.paragraphs:
                text += paragraph.text + "\n"
            return text
        except Exception as e:
            logger.error(f"DOCX extraction failed: {e}")
            raise
    
    def _extract_xlsx_text(self, content: bytes) -> str:
        """Extract text from XLSX"""
        try:
            workbook = openpyxl.load_workbook(io.BytesIO(content))
            text = ""
            for sheet_name in workbook.sheetnames:
                worksheet = workbook[sheet_name]
                text += f"Sheet: {sheet_name}\n"
                for row in worksheet.iter_rows(values_only=True):
                    row_text = " ".join(str(cell) for cell in row if cell)
                    if row_text.strip():
                        text += row_text + "\n"
                text += "\n"
            return text
        except Exception as e:
            logger.error(f"XLSX extraction failed: {e}")
            raise
    
    def _extract_pptx_text(self, content: bytes) -> str:
        """Extract text from PPTX"""
        try:
            prs = Presentation(io.BytesIO(content))
            text = ""
            for slide_num, slide in enumerate(prs.slides, 1):
                text += f"Slide {slide_num}:\n"
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        text += shape.text + "\n"
                text += "\n"
            return text
        except Exception as e:
            logger.error(f"PPTX extraction failed: {e}")
            raise
    
    async def _store_document_embeddings(self, text: str, document_id: str, filename: str) -> int:
        """Store document embeddings in vector database"""
        try:
            # Split text into chunks
            chunks = self._split_text_into_chunks(text, chunk_size=1000, overlap=200)
            
            embeddings_count = 0
            for i, chunk in enumerate(chunks):
                if len(chunk.strip()) < 50:  # Skip very short chunks
                    continue
                
                # Create embeddings
                embedding = self.embedding_model.encode(chunk)
                
                # Store in ChromaDB
                self.collection.add(
                    embeddings=[embedding.tolist()],
                    documents=[chunk],
                    metadatas=[{
                        'document_id': document_id,
                        'filename': filename,
                        'chunk_index': i,
                        'type': 'document',
                        'timestamp': datetime.now().isoformat()
                    }],
                    ids=[f"doc_{document_id}_{i}"]
                )
                embeddings_count += 1
            
            return embeddings_count
            
        except Exception as e:
            logger.error(f"Failed to store document embeddings: {e}")
            return 0
    
    def _split_text_into_chunks(self, text: str, chunk_size: int = 1000, overlap: int = 200) -> List[str]:
        """Split text into overlapping chunks"""
        chunks = []
        start = 0
        
        while start < len(text):
            end = start + chunk_size
            chunk = text[start:end]
            
            # Try to break at sentence boundary
            if end < len(text):
                last_period = chunk.rfind('.')
                last_newline = chunk.rfind('\n')
                break_point = max(last_period, last_newline)
                
                if break_point > chunk_size * 0.7:  # Only break if we're not too early
                    chunk = chunk[:break_point + 1]
                    end = start + break_point + 1
            
            chunks.append(chunk.strip())
            start = end - overlap
        
        return chunks
    
    async def search_documents(self, query: str, limit: int = 5) -> List[Dict[str, Any]]:
        """Search documents using semantic similarity"""
        try:
            # Create query embedding
            query_embedding = self.embedding_model.encode(query)
            
            # Search in ChromaDB
            results = self.collection.query(
                query_embeddings=[query_embedding.tolist()],
                n_results=limit,
                include=['documents', 'metadatas', 'distances']
            )
            
            return [
                {
                    'content': doc,
                    'metadata': meta,
                    'relevance_score': 1 - distance
                }
                for doc, meta, distance in zip(
                    results['documents'][0],
                    results['metadatas'][0],
                    results['distances'][0]
                )
            ]
            
        except Exception as e:
            logger.error(f"Document search failed: {e}")
            return [] 